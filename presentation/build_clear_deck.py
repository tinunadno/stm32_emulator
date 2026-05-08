#!/usr/bin/env python3
"""
Краткая «максимально понятная» презентация (≈40–45 слайдов):
крупный текст, на каждом слайде пояснение «что это / зачем», только суть проекта.

  python3 presentation/build_clear_deck.py

Нужны: python-pptx, cairosvg, matplotlib (как у build_full_md_deck.py).
"""

from __future__ import annotations

import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
HERE = Path(__file__).resolve().parent
OUT = HERE / "Защита_понятно_кратко.pptx"

C_BG = RGBColor(0xFA, 0xFB, 0xFD)
C_TITLE = RGBColor(0x0D, 0x2B, 0x5E)
C_BODY = RGBColor(0x1A, 0x1A, 0x1A)
C_ACCENT = RGBColor(0x0B, 0x5C, 0x4F)
C_MUTED = RGBColor(0x55, 0x55, 0x55)
C_COVER = RGBColor(0x12, 0x2A, 0x4A)


def load_image_generator():
    import sys

    path = HERE / "build_full_md_deck.py"
    name = "_stm32_deck_images"
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    assert spec.loader
    spec.loader.exec_module(mod)
    return mod.generate_all_images


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, rgb):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb


def slide_cover(prs, title: str, subtitle: str):
    s = blank(prs)
    bg(s, C_COVER)
    t = s.shapes.add_textbox(Inches(0.55), Inches(2.1), Inches(8.9), Inches(1.4))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xC5, 0xD8, 0xF0)
    p2.space_before = Pt(14)


def slide_explained(
    prs,
    title: str,
    explanation: str,
    bullets: list[str] | None = None,
    footer: str = "",
):
    """Один слайд: заголовок + обязательное пояснение + немного маркеров."""
    s = blank(prs)
    bg(s, C_BG)
    y = 0.38
    tb = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9.0), Inches(0.75))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = C_TITLE
    y += 0.82

    ex = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9.0), Inches(1.55))
    ef = ex.text_frame
    ef.word_wrap = True
    ep = ef.paragraphs[0]
    ep.text = explanation
    ep.font.size = Pt(18)
    ep.font.color.rgb = C_BODY
    ep.line_spacing = 1.25
    y += 1.65

    if bullets:
        bx = s.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9.0), Inches(4.5))
        bf = bx.text_frame
        bf.word_wrap = True
        for i, b in enumerate(bullets[:6]):
            para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            para.text = f"• {b}"
            para.font.size = Pt(17)
            para.font.color.rgb = C_ACCENT
            para.space_after = Pt(6)

    if footer:
        fb = s.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9.0), Inches(0.5))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(12)
        fp.font.color.rgb = C_MUTED


def slide_figure(
    prs,
    title: str,
    before: str,
    png: Path,
    after: str,
    pic_height=Inches(3.35),
):
    s = blank(prs)
    bg(s, C_BG)
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.65))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(24)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = C_TITLE

    b1 = s.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9.0), Inches(0.95))
    p1 = b1.text_frame.paragraphs[0]
    p1.text = before
    p1.font.size = Pt(15)
    p1.font.color.rgb = C_BODY

    if png.is_file():
        s.shapes.add_picture(str(png), Inches(0.45), Inches(1.95), height=pic_height)

    b2 = s.shapes.add_textbox(Inches(0.5), Inches(5.45), Inches(9.0), Inches(1.35))
    p2 = b2.text_frame.paragraphs[0]
    p2.text = after
    p2.font.size = Pt(15)
    p2.font.color.rgb = C_BODY


def slide_code(prs, title: str, explain: str, lines: list[str]):
    s = blank(prs)
    bg(s, RGBColor(0xF6, 0xF8, 0xFA))
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.6))
    t.text_frame.paragraphs[0].text = title
    t.text_frame.paragraphs[0].font.size = Pt(22)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = C_TITLE
    e = s.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9.0), Inches(0.7))
    e.text_frame.paragraphs[0].text = explain
    e.text_frame.paragraphs[0].font.size = Pt(14)
    e.text_frame.paragraphs[0].font.color.rgb = C_BODY
    code = "\n".join(lines)
    c = s.shapes.add_textbox(Inches(0.45), Inches(1.65), Inches(9.1), Inches(5.0))
    cp = c.text_frame.paragraphs[0]
    cp.text = code
    cp.font.name = "Courier New"
    cp.font.size = Pt(11)
    cp.font.color.rgb = RGBColor(0x11, 0x11, 0x11)


def build():
    gen = load_image_generator()
    img = gen()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- 1–8: вводная ---
    slide_cover(
        prs,
        "Симулятор STM32F103C8T6",
        "Курс «Программирование встроенных систем» · защита проекта",
    )
    slide_explained(
        prs,
        "Один слайд — суть проекта",
        "Мы сделали программу на компьютере, которая ведёт себя как микроконтроллер STM32F103: "
        "выполняет вашу прошивку, обрабатывает прерывания, таймер и UART. Плата не обязательна.",
        bullets=["Язык: C11", "Запуск: свой бинарник или пример из examples/", "Есть CLI и удалённая отладка GDB"],
    )
    slide_explained(
        prs,
        "Зачем это нужно",
        "На лекциях и лабораторных важно видеть, как работает процессор и периферия. "
        "Симулятор даёт пошаговый запуск, лог и точки останова без риска «сжечь» железо.",
        bullets=["Учёба и отладка алгоритмов", "Повторяемый сценарий на любой машине", "Можно сдавать работу с демонстрацией на ноутбуке"],
    )
    slide_explained(
        prs,
        "Что такое STM32F103 / «Blue Pill»",
        "Это недорогой микроконтроллер с ядром ARM Cortex-M3. В нём есть Flash (программа), "
        "SRAM (данные), таймеры, UART и система прерываний. Наш проект повторяет эту модель в программе.",
        bullets=["Cortex-M3 — 32-битное ядро, режим Thumb/Thumb-2", "Память и регистры периферии — по адресам, как в документации"],
    )
    slide_explained(
        prs,
        "Что значит «симулятор» в этом проекте",
        "Это не эмуляция «в реальном времени как железо», а детерминированная модель: "
        "один поток, фиксированный порядок шагов. Каждый шаг = обновление периферии, исполнение инструкции, проверка IRQ и breakpoints.",
        bullets=["Предсказуемо: одинаковый вход → одинаковый результат", "Удобно для тестов и отладки"],
    )
    slide_explained(
        prs,
        "Что требовал курс (кратко)",
        "В ТЗ (README.md) заложены: ядро Cortex-M3, память, NVIC, таймер, UART, загрузка прошивки, тесты. "
        "Отчёт и презентация — отдельные артефакты сдачи.",
        bullets=["Групповой проект до 6 человек", "Модульные и интеграционные тесты", "Демонстрация работы"],
        footer="Источник требований: README.md",
    )
    slide_explained(
        prs,
        "Какие документы к чему относятся",
        "Чтобы не путаться на защите: README.md — учебное ТЗ; IMPLEMENTATION.md — как устроен код; "
        "development_plan/description.md — формальные требования (часть текста про C++, фактически пишем на C11); "
        "docs/gateway.md — идея будущей веб-оболочки, не текущая программа.",
        bullets=["Всегда опирайтесь на IMPLEMENTATION.md для фактов о коде", "gateway.md — только перспектива"],
    )
    slide_explained(
        prs,
        "Что уже есть в коде (факты)",
        "В каталоге src/ лежит симулятор на C11: ядро, шина, Flash и SRAM, NVIC, таймер TIM2, UART USART1, GPIO, RCC, SysTick, "
        "точки останова, консоль и GDB-сервер. Это и есть результат команды; веб-очередь из gateway.md в этот репозиторий не входит.",
        bullets=["Сборка: cd src && make", "Запуск: ./stm32sim [файл.bin] [--gdb]", "Проверки: make test", "Пример прошивки: examples/"],
        footer="См. IMPLEMENTATION.md",
    )

    slide_code(
        prs,
        "Как собрать и запустить (терминал)",
        "Сначала собираем симулятор. Бинарник прошивки — обычно .bin для загрузки во Flash.",
        [
            "cd src",
            "make              # stm32sim",
            "./stm32sim                    # интерактивный CLI",
            "./stm32sim firmware.bin       # сразу с прошивкой",
            "./stm32sim firmware.bin --gdb # сервер GDB (порт по умолчанию см. gdb_stub)",
        ],
    )
    slide_explained(
        prs,
        "Распределение ролей (как в ТЗ)",
        "Так в README.md предлагают разделить работу шестерым. На слайде — ориентир: кто за какую область отвечает.",
        bullets=[
            "Архитектор / ядро CPU",
            "Память и загрузка бинарника",
            "NVIC и прерывания",
            "Таймер",
            "UART и консоль",
            "Тесты, сборка, документация",
        ],
        footer="README.md, § 8",
    )
    slide_explained(
        prs,
        "Этапы работ (логика семестра)",
        "Типичный порядок: спроектировать → сделать CPU и память → NVIC → таймер → UART → связать → тесты → отчёт и слайды.",
        bullets=["Не обязательно строго по неделям — важен итог: работающий симулятор и проверки"],
        footer="README.md, § 9",
    )
    slide_explained(
        prs,
        "Два режима: консоль и отладчик GDB",
        "Без флага --gdb вы работаете в текстовом интерфейсе: вводите команды сами. С флагом --gdb симулятор открывает сетевой порт: "
        "к нему подключается arm-none-eabi-gdb (как к удалённой плате). Так удобнее показывать breakpoints и обход кода на защите.",
        bullets=["CLI — быстро показать память и регистры", "GDB — если привыкли к отладке из IDE или терминала"],
    )
    slide_explained(
        prs,
        "Консольные команды (самое нужное)",
        "После запуска ./stm32sim без --gdb вы попадаете в CLI. Ниже команды, которые чаще всего показывают на защите.",
        bullets=[
            "help — список команд",
            "load путь — загрузить .bin во Flash",
            "step [N] — выполнить N инструкций",
            "run / stop — бег / пауза до breakpoint",
            "reg / mem адрес — состояние",
            "break адрес / delete адрес — точки останова",
            "uart символ — подать байт на приёмник",
        ],
        footer="Полная таблица: IMPLEMENTATION.md",
    )

    # --- Рисунки (после того, как объяснили словами) ---
    slide_figure(
        prs,
        "Рисунок 1. Общая архитектура симулятора",
        "Ниже — схема из docs/architecture.svg. Она показывает, как ядро, память, шина и периферия связаны между собой.",
        img["architecture"],
        "Читать так: ядро исполняет инструкции и ходит в память через шину. "
        "NVIC решает, когда прервать выполнение. Периферия (таймер, UART…) поднимает запросы к NVIC.",
        pic_height=Inches(3.15),
    )
    slide_figure(
        prs,
        "Рисунок 2. Папки исходников (упрощённо)",
        "Это не все файлы, а логические блоки: кто за что отвечает в каталоге src/.",
        img["modules"],
        "core — процессор; memory — Flash/RAM; bus — маршрутизация адресов; nvic — прерывания; "
        "peripherals — устройства; simulator — главный цикл; ui — консоль; gdb_stub — отладка по сети.",
        pic_height=Inches(3.05),
    )
    slide_figure(
        prs,
        "Рисунок 3. Карта памяти (адреса)",
        "Таблица повторяет раздел IMPLEMENTATION.md. Адреса нужны, чтобы прошивка «видела» те же регистры, что и на кристалле.",
        img["memory_map"],
        "Flash — код; SRAM — стек и данные; по адресам 0x4000… лежат регистры TIM2 и USART1. "
        "Если адрес не распознан — шина сообщает об ошибке доступа.",
        pic_height=Inches(2.85),
    )
    slide_figure(
        prs,
        "Рисунок 4. Один «такт» симуляции",
        "Так выглядит строгий порядок из IMPLEMENTATION.md. Его нельзя переставлять — от этого зависит детерминизм.",
        img["sim_cycle"],
        "Сначала тикают все периферии (таймеры могут досчитать и запросить IRQ). Потом выполняется одна инструкция процессора. "
        "Затем NVIC решает, нужно ли войти в прерывание. В конце проверяются точки останова.",
        pic_height=Inches(2.55),
    )

    # --- 19–24: модули простым языком ---
    slide_explained(
        prs,
        "Ядро (core/)",
        "Декодирует 16-битные инструкции Thumb, обновляет регистры R0–R15 и флаги, обращается к шине за данными и инструкциями. "
        "Набор инструкций покрывает типичные операции MOV, LDR/STR, арифметику, ветвления, PUSH/POP и др. (см. таблицу в IMPLEMENTATION.md).",
        bullets=["Одна инструкция за вызов core_step (в рамках общего цикла)", "Поддержка входа/выхода из исключений NVIC"],
    )
    slide_explained(
        prs,
        "Память (memory/) и шина (bus/)",
        "Память хранит образ Flash и SRAM. Шина по адресу решает: это обращение к RAM, Flash или к регистру периферии. "
        "Так мы не даём ядру «знать» про устройства напрямую — только адрес.",
        bullets=["Little-endian, как на ARM", "Регистрация регионов — расширяемо"],
    )
    slide_explained(
        prs,
        "Старт прошивки: откуда берётся PC и SP",
        "После сброса процессор читает первые слова из Flash: указатель стека и адрес Reset_Handler. "
        "Если это не совпадает с вашим linker script, прошивка «не оживёт» — на защите стоит сказать, что бинарник собран под ту же карту памяти, что в симуляторе.",
        bullets=["SP ← слово по 0x08000000", "PC ← слово по 0x08000004", "Дальше идёт ваш код"],
        footer="См. гарантии в description.md / IMPLEMENTATION.md",
    )

    slide_explained(
        prs,
        "NVIC — что происходит при прерывании (без формул)",
        "Когда IRQ разрешён и приоритет позволяет, процессор прерывает обычный код, кладёт часть регистров в стек и перескакивает на обработчик из таблицы векторов. "
        "После обработчика — возврат по специальному значению LR (EXC_RETURN). В проекте заложено до 43 линий запросов — как в спецификации симулятора.",
        bullets=["Контекст сохраняется по правилам Cortex-M (в модели упрощённо, но предсказуемо)", "Приоритет: меньше число — важнее прерывание"],
    )

    slide_explained(
        prs,
        "Периферия: TIM2 и UART",
        "TIM2 считает такты с прескалером и может генерировать IRQ по переполнению — типичный сценарий «тик раз в N циклов». "
        "USART1 принимает и передаёт байты; прерывания помогают прошивке реагировать на приём и окончание передачи.",
        bullets=["Регистры по адресам из карты памяти", "Связь с NVIC через линии IRQ"],
    )
    slide_figure(
        prs,
        "Рисунок 5. UART на временной диаграмме",
        "Файл uart_diagram.svg — пример визуализации: когда какой байт ушел по TX. Полезно объяснять работу прошивки на защите.",
        img["uart"],
        "Каждый «столбик» — символ на линии TX. Внизу — условная шкала тактов симуляции. "
        "Это наглядно, если преподаватель спросит: «покажите, что UART реально шевелится».",
        pic_height=Inches(2.65),
    )

    slide_explained(
        prs,
        "Отладчик и GDB-stub",
        "Модуль debugger/ хранит список адресов останова. gdb_stub/ поднимает TCP-сервер протокола GDB Remote Serial Protocol — "
        "к нему подключается arm-none-eabi-gdb с командой target remote …",
        bullets=["Breakpoints останавливают цикл симуляции", "GDB — для привычной отладки студентам embedded"],
    )

    # --- 25–28: тесты и пример ---
    slide_code(
        prs,
        "Тесты — как убедиться, что всё не сломано",
        "Команда make test собирает test_runner и гоняет проверки по модулям. Это ответ на вопрос «где доказательство корректности».",
        [
            "cd src",
            "make test",
            "",
            "# В каталоге tests/: core, memory, bus, nvic,",
            "# timer, uart, gpio, интеграция и др.",
        ],
    )
    slide_explained(
        prs,
        "Пример прошивки (examples/)",
        "Там Makefile под arm-none-eabi-gcc и простой код старта. Вы собираете firmware.bin и подаёте его в stm32sim — "
        "так показывают работу «как на лабораторной».",
        bullets=["Свой linker script и вектор сброса", "Можно менять код и смотреть вывод UART в симуляторе"],
    )

    slide_explained(
        prs,
        "Как показать работу на защите (сценарий на 2 минуты)",
        "Краткий сценарий, чтобы не теряться: (1) собрать симулятор и сказать, что это за бинарник; (2) запустить с examples/ или своим .bin; "
        "(3) выполнить несколько step или run, показать reg/mem или UART; (4) при желании — подключить GDB и поставить breakpoint.",
        bullets=["Заранее проверьте команды на своём ноутбуке", "Если что-то не собирается — честно скажите и покажите тесты make test"],
    )

    # --- 29–32: проектирование и будущее ---
    slide_explained(
        prs,
        "Как проект расширяют без хаоса",
        "Новая периферия подключается через таблицу функций (vtable): read/write/tick/reset. "
        "Новая инструкция — новая строка в таблице декодера. Новая команда CLI — запись в массив команд.",
        bullets=["Идея Open/Closed: меньше правок в старом коде", "Подробно: IMPLEMENTATION.md"],
    )
    slide_explained(
        prs,
        "Про формальную спецификацию в development_plan/",
        "В description.md текст иногда описывает C++17 — это черновик для обсуждения. Фактическая реализация в репозитории — на C11, как в IMPLEMENTATION.md.",
        bullets=["На защите важно говорить: «реализация = C11»", "Спека = ориентир по поведению, не по синтаксису"],
    )
    slide_figure(
        prs,
        "Рисунок 6. Идея онлайн-лаборатории (только документ)",
        "docs/gateway.md — это про будущее: очередь заданий, воркеры, KeyDB, браузер и GDB. В текущем репозитории этого кода нет.",
        img["gateway"],
        "Зачем слайд: показать дорожную карту. Сейчас продукт — локальный симулятор; распределённая оболочка — отдельный этап, если команда продолжит.",
        pic_height=Inches(2.75),
    )

    slide_explained(
        prs,
        "Почему симулятор «детерминированный»",
        "Мы специально не используем потоки и «живые» таймеры ОС: каждый запуск с теми же входными данными даёт тот же сценарий. "
        "Это упрощает отладку и позволяет тестам стабильно ловить ошибки.",
        bullets=["Нет гонок потоков", "Нет случайных задержек", "Поведение воспроизводимо на любой машине"],
    )
    slide_explained(
        prs,
        "Ошибки в коде: как они сообщаются",
        "Функции возвращают код Status вместо исключений C++. Так проще сочетать чистый C и предсказуемые пути выполнения.",
        bullets=["Успех / ошибка памяти / неверный адрес / нереализованная инструкция — всё явно", "Вызовы проверяют код и не «молчат»"],
    )
    slide_explained(
        prs,
        "Прошивка: .bin и связь с линкером",
        "Симулятор загружает готовый бинарник в модель Flash. Адреса в бинарнике должны совпадать с картой памяти (0x08000000…). "
        "Если линкер кладёт код не туда, симулятор не «починит» это сам — это нормально для учебного проекта.",
        bullets=["Проверьте vector table и entry point перед демонстрацией", "Пример правильной связки — examples/"],
    )

    slide_explained(
        prs,
        "Три термина, которые любят спросить",
        "MMIO (memory-mapped I/O) — когда регистр устройства читается как ячейка памяти по фиксированному адресу. "
        "IRQ — запрос прерывания от периферии к NVIC. GDB RSP — текстовый протокол между gdb и удалённым «заглушкой» (наш gdb_stub).",
        bullets=["Все три есть в проекте в явном виде", "Не нужно заучивать стандарт — достаточно объяснить смысл"],
    )
    slide_explained(
        prs,
        "Что дают интеграционные тесты",
        "Юнит-тесты проверяют кусочки; интеграционный сценарий убеждает, что CPU, память, NVIC и таймер согласованы — например, что IRQ действительно вызывает обработчик.",
        bullets=["См. test_integration.c и Makefile", "На защите: «зелёный make test» — сильный аргумент"],
    )

    # --- 33–end ---
    slide_explained(
        prs,
        "Ограничения и честный статус",
        "Версия в version.txt — ранняя. Не заявляем полное соответствие всем тонкостям кремния STM32: это учебная модель с фокусом на архитектуру и практику.",
        bullets=["Нет полного набора периферии чипа", "Нет CI в репозитории — тесты локальные", "Онлайн-сервис из gateway.md не реализован"],
    )
    slide_explained(
        prs,
        "Итог для проверяющего",
        "Есть работающий симулятор с документацией, тестами, CLI, GDB и примером прошивки. Архитектура объясняется схемами и картой памяти. "
        "План развития (онлайн) отделён от текущего кода.",
        bullets=["Собирается и запускается стандартными командами", "Понятная модульная структура", "Наглядные рисунки в презентации"],
    )
    slide_explained(
        prs,
        "Команда (заполнить)",
        "Укажите роли: кто ядро, память/NVIC, периферия, тесты, документация, интеграция.",
        bullets=[
            "1. ФИО, группа — _______________",
            "2. ФИО, группа — _______________",
            "3. ФИО, группа — _______________",
            "4. ФИО, группа — _______________",
            "5. ФИО, группа — _______________",
            "6. ФИО, группа — _______________",
        ],
    )
    slide_cover(prs, "Спасибо за внимание", "Готовы ответить на вопросы")

    prs.save(str(OUT))
    n = len(prs.slides)
    print(f"Saved {OUT} ({n} slides)")


if __name__ == "__main__":
    build()
