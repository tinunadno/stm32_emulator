#!/usr/bin/env python3
"""Генерация PPTX для защиты проекта stm32_emulator (курс ПрВС)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Палитра (тёмно-синий / бирюза — тема «встраиваемые системы»)
C_BG_DARK = RGBColor(0x21, 0x29, 0x5C)
C_PRIMARY = RGBColor(0x06, 0x5A, 0x82)
C_ACCENT = RGBColor(0x02, 0xC3, 0x9A)
C_TEXT = RGBColor(0x2A, 0x2A, 0x2A)
C_MUTED = RGBColor(0x55, 0x55, 0x55)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL = RGBColor(0xF0, 0xF4, 0xF8)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def slide_title_dark(prs, title: str, subtitle: str):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_BG_DARK)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)
    p2.space_before = Pt(12)

    foot = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.6))
    ft = foot.text_frame
    fp = ft.paragraphs[0]
    fp.text = "ИТМО · Программирование встроенных систем · 2026"
    fp.font.size = Pt(12)
    fp.font.color.rgb = RGBColor(0xAA, 0xB8, 0xCC)


def slide_section_bar(prs, section: str):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(3.0), Inches(8.8), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = section
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_WHITE


def slide_bullets(prs, title: str, bullets: list[str], footnote: str | None = None):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_PANEL)

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.9), Inches(0.85))
    tfp = tb.text_frame
    tp = tfp.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = C_PRIMARY

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(8.9), Inches(5.5))
    bf = body.text_frame
    bf.word_wrap = True
    for i, line in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = C_TEXT
        para.space_after = Pt(8)
        para.level = 0

    if footnote:
        fn = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(8.9), Inches(0.55))
        ff = fn.text_frame
        fp = ff.paragraphs[0]
        fp.text = footnote
        fp.font.size = Pt(11)
        fp.font.color.rgb = C_MUTED
        fp.font.italic = True


def slide_two_col(prs, title: str, left_title: str, left_lines: list[str], right_title: str, right_lines: list[str]):
    slide = _blank_slide(prs)
    _fill_bg(slide, C_PANEL)

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(8.9), Inches(0.8))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = C_PRIMARY

    w = Inches(4.2)
    y = Inches(1.25)
    for x, st, lines in [
        (Inches(0.55), left_title, left_lines),
        (Inches(5.0), right_title, right_lines),
    ]:
        box = slide.shapes.add_textbox(x, y, w, Inches(5.4))
        tf = box.text_frame
        tf.word_wrap = True
        h = tf.paragraphs[0]
        h.text = st
        h.font.size = Pt(18)
        h.font.bold = True
        h.font.color.rgb = C_BG_DARK
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = C_TEXT
            p.space_after = Pt(6)


def main():
    root = Path(__file__).resolve().parents[1]
    out = Path(__file__).resolve().parent / "Защита_STM32_симулятор.pptx"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title_dark(
        prs,
        "Симулятор STM32F103C8T6",
        "ARM Cortex-M3 · периферия · GDB · тестирование\nКурс «Программирование встроенных систем»",
    )

    slide_bullets(
        prs,
        "Цель проекта",
        [
            "Программный симулятор МК STM32F103C8T6 (класс Blue Pill) на C11.",
            "Исполнение прошивки без физической платы: учёба, лабораторные, отладка логики.",
            "Соответствие групповому ТЗ в README.md (ядро, память, NVIC, таймер, UART, тесты).",
        ],
        footnote="Источник: README.md, IMPLEMENTATION.md",
    )

    slide_two_col(
        prs,
        "Что реализовано в репозитории",
        "Ядро и платформа",
        [
            "Cortex-M3: Thumb / Thumb-2 (core/)",
            "Flash 64 KB, SRAM 20 KB (memory/)",
            "Системная шина (bus/)",
            "NVIC, 43 линии (nvic/)",
        ],
        "Периферия и инструменты",
        [
            "TIM2, SysTick, RCC, GPIO, USART1",
            "Точки останова (debugger/)",
            "CLI (ui/)",
            "GDB Remote Serial Protocol (gdb_stub/)",
        ],
    )

    slide_bullets(
        prs,
        "Архитектура (модульность)",
        [
            "Оркестратор симуляции: simulator/ — цикл tick / step / прерывания.",
            "Расширение периферии через vtable Peripheral (Open/Closed в C).",
            "Внутренняя очередь событий events/ — планирование по циклам (не HTTP SSE).",
        ],
        footnote="Подробно: IMPLEMENTATION.md",
    )

    slide_bullets(
        prs,
        "Отладка и проверка",
        [
            "Запуск: cd src && make → ./stm32sim [firmware.bin], опция GDB.",
            "Автотесты: make test — NVIC, память, шина, таймеры, UART, ядро, интеграция.",
            "Пример прошивки: examples/ (arm-none-eabi-gcc) → бинарник для симулятора.",
        ],
        footnote="Тесты: src/tests/, пример: examples/Makefile",
    )

    slide_bullets(
        prs,
        "Документация и перспектива",
        [
            "ТЗ и требования к курсовому проекту — README.md.",
            "Концепция онлайн-лаборатории (очередь заданий, API, масштабирование) — docs/gateway.md.",
            "Схема архитектуры — docs/architecture.svg.",
            "Важно для защиты: распределённый стек из gateway.md в коде не реализован — это этап развития / сопутствующий дизайн.",
        ],
        footnote="Честно отделяем «сделано в C» от «спроектировано в документации».",
    )

    slide_bullets(
        prs,
        "Команда проекта (до 6 чел.)",
        [
            "ФИО, группа — ___________________",
            "ФИО, группа — ___________________",
            "ФИО, группа — ___________________",
            "ФИО, группа — ___________________",
            "ФИО, группа — ___________________",
            "ФИО, группа — ___________________",
        ],
        footnote="Заполните перед защитой или замените на реальные роли (ядро, периферия, тесты, документация).",
    )

    slide_bullets(
        prs,
        "Результат",
        [
            "Детерминированный однопоточный симулятор с документированной структурой.",
            "Версия (version.txt): 0.0.1 — база для расширения периферии и интеграций.",
            "Скрипт пересборки презентации: presentation/build_defense_deck.py",
        ],
    )

    slide_section_bar(prs, "Спасибо за внимание")
    slide = prs.slides[-1]
    box = slide.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(8.8), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "Вопросы?"
    p.font.size = Pt(24)
    p.font.color.rgb = C_WHITE

    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
