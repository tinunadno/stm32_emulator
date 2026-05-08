#!/usr/bin/env python3
"""
Сборка PPTX из всех проектных .md + вставка сгенерированных и исходных изображений.

Использование (из корня репозитория):
  python3 presentation/build_full_md_deck.py

Выход:
  presentation/assets_generated/*.png
  presentation/Проект_полный_MD_и_рисунки.pptx
"""

from __future__ import annotations

import re
import textwrap
from dataclasses import dataclass
from pathlib import Path

import cairosvg
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
PRESENTATION = Path(__file__).resolve().parent
ASSETS = PRESENTATION / "assets_generated"
OUT_PPTX = PRESENTATION / "Проект_полный_MD_и_рисунки.pptx"

MD_SOURCES: list[tuple[Path, str]] = [
    (ROOT / "README.md", "README.md — ТЗ курса"),
    (ROOT / "IMPLEMENTATION.md", "IMPLEMENTATION.md — реализация"),
    (ROOT / "development_plan" / "description.md", "development_plan/description.md — спецификация"),
    (ROOT / "docs" / "gateway.md", "docs/gateway.md — онлайн-лаборатория (дизайн)"),
]

C_BG = RGBColor(0xF5, 0xF7, 0xFA)
C_TITLE = RGBColor(0x0D, 0x47, 0xA1)
C_TEXT = RGBColor(0x22, 0x22, 0x22)
C_MUTED = RGBColor(0x66, 0x66, 0x66)
C_COVER = RGBColor(0x15, 0x2C, 0x52)

MAX_BULLET_CHARS = 220
MAX_BULLETS_PER_SLIDE = 11
MAX_CODE_LINES = 20


def ensure_assets() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)


def svg_to_png(svg: Path, png: Path, scale: float = 2.0) -> None:
    if not svg.is_file():
        return
    cairosvg.svg2png(url=str(svg), write_to=str(png), scale=scale)


def gen_memory_map_png(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 3.2))
    ax.axis("off")
    rows = [
        ["Flash (alias)", "0x00000000", "64 KB", "Векторная таблица"],
        ["Flash", "0x08000000", "64 KB", "Прошивка"],
        ["SRAM", "0x20000000", "20 KB", "ОЗУ"],
        ["TIM2", "0x40000000", "1 KB", "Регистры таймера"],
        ["USART1", "0x40013800", "1 KB", "Регистры UART"],
    ]
    col = ["Регион", "База", "Размер", "Описание"]
    t = ax.table(
        cellText=rows,
        colLabels=col,
        loc="center",
        cellLoc="left",
    )
    t.auto_set_font_size(False)
    t.set_fontsize(9)
    t.scale(1.05, 1.8)
    fig.suptitle("Карта памяти (IMPLEMENTATION.md)", fontsize=12, fontweight="bold")
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close()


def gen_simulation_cycle_png(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 2.8))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 2)
    ax.axis("off")
    steps = [
        (0.3, "1. tick_peripherals()\nпериферия"),
        (2.5, "2. core_step()\nинструкция"),
        (4.7, "3. check_interrupts()\nNVIC"),
        (6.9, "4. check_breakpoints()\nотладчик"),
    ]
    for x, txt in steps:
        box = FancyBboxPatch(
            (x, 0.55), 1.7, 1.0, boxstyle="round,pad=0.05",
            linewidth=1.5, edgecolor="#1565C0", facecolor="#E3F2FD",
        )
        ax.add_patch(box)
        ax.text(x + 0.85, 1.05, txt, ha="center", va="center", fontsize=8.5)
    for i in range(len(steps) - 1):
        ax.annotate(
            "",
            xy=(steps[i + 1][0] - 0.05, 1.05),
            xytext=(steps[i][0] + 1.75, 1.05),
            arrowprops=dict(arrowstyle="->", color="#333", lw=1.5),
        )
    fig.suptitle("Цикл симуляции (фиксированный порядок)", fontsize=11, fontweight="bold")
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close()


def gen_module_graph_png(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 6))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 7)
    ax.axis("off")

    def box(xy, w, h, label, fc, ec):
        b = FancyBboxPatch(
            xy, w, h, boxstyle="round,pad=0.03",
            linewidth=1.2, edgecolor=ec, facecolor=fc,
        )
        ax.add_patch(b)
        ax.text(xy[0] + w / 2, xy[1] + h / 2, label, ha="center", va="center", fontsize=8.5)

    box((4.2, 5.0), 2.6, 0.75, "Cortex-M3 core/", "#E8F5E9", "#2E7D32")
    box((0.4, 3.0), 2.0, 0.7, "memory/", "#FFF3E0", "#EF6C00")
    box((4.5, 3.0), 2.0, 0.7, "bus/", "#E3F2FD", "#1565C0")
    box((8.2, 3.0), 2.0, 0.7, "nvic/", "#F3E5F5", "#6A1B9A")
    box((0.3, 1.0), 1.6, 0.65, "timer/", "#ECEFF1", "#455A64")
    box((2.3, 1.0), 1.6, 0.65, "uart/\ngpio/rcc…", "#ECEFF1", "#455A64")
    box((4.3, 1.0), 1.6, 0.65, "simulator/", "#FFFDE7", "#F9A825")
    box((6.3, 1.0), 1.6, 0.65, "debugger/", "#FFEBEE", "#C62828")
    box((8.3, 1.0), 1.6, 0.65, "ui/ gdb_stub/", "#E0F7FA", "#00838F")

    arrows = [
        ((5.0, 5.0), (5.5, 3.7)),
        ((5.5, 3.35), (5.5, 1.65)),
        ((4.2, 5.35), (1.4, 3.7)),
        ((6.8, 5.35), (9.2, 3.7)),
        ((5.1, 3.35), (1.1, 1.65)),
        ((5.9, 3.35), (3.1, 1.65)),
    ]
    for (x1, y1), (x2, y2) in arrows:
        ax.add_patch(
            FancyArrowPatch(
                (x1, y1), (x2, y2),
                arrowstyle="-|>", mutation_scale=12, color="#555", linewidth=1,
            )
        )
    fig.suptitle("Модули src/ (упрощённая схема)", fontsize=12, fontweight="bold")
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close()


def gen_gateway_overview_png(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12, 6.5))
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 8)
    ax.axis("off")

    def b(xy, w, h, t, c):
        patch = FancyBboxPatch(
            xy, w, h, boxstyle="round,pad=0.04",
            linewidth=1.2, edgecolor="#0D47A1", facecolor=c,
        )
        ax.add_patch(patch)
        ax.text(xy[0] + w / 2, xy[1] + h / 2, t, ha="center", va="center", fontsize=8)

    b((0.3, 6.2), 1.6, 1.0, "Клиент\n(браузер /\nGDB)", "#BBDEFB")
    b((2.5, 6.0), 3.2, 1.4, "API Gateway (Go)\nREST + SSE", "#90CAF9")
    b((6.5, 6.0), 2.2, 1.4, "PostgreSQL\nаудит", "#C8E6C9")
    b((2.8, 3.5), 5.5, 1.2, "KeyDB — очереди, job:{id}, Pub/Sub событий", "#FFE082")
    b((0.5, 0.8), 2.4, 1.0, "Worker 1\nсимуляция +\nGDB порт", "#FFCCBC")
    b((3.3, 0.8), 2.4, 1.0, "Worker 2", "#FFCCBC")
    b((6.1, 0.8), 2.4, 1.0, "Worker N", "#FFCCBC")
    b((9.0, 0.8), 2.0, 1.0, "GDB клиенты\nTCP напрямую", "#D1C4E9")

    flows = [
        ((1.1, 6.2), (2.5, 6.7)),
        ((5.7, 6.0), (6.5, 6.7)),
        ((4.5, 6.0), (4.55, 4.7)),
        ((4.0, 3.5), (1.7, 1.8)),
        ((5.5, 3.5), (4.5, 1.8)),
        ((6.5, 3.5), (7.3, 1.8)),
        ((8.5, 1.3), (9.0, 1.3)),
    ]
    for a, b_ in flows:
        ax.add_patch(FancyArrowPatch(a, b_, arrowstyle="-|>", mutation_scale=14, color="#444", lw=1.1))
    fig.suptitle("Онлайн-лаборатория (docs/gateway.md) — обзор", fontsize=12, fontweight="bold")
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor="white")
    plt.close()


def generate_all_images() -> dict[str, Path]:
    ensure_assets()
    mapping: dict[str, Path] = {}
    arch = ASSETS / "architecture.png"
    uart = ASSETS / "uart_diagram.png"
    svg_to_png(ROOT / "docs" / "architecture.svg", arch)
    svg_to_png(ROOT / "uart_diagram.svg", uart)
    mapping["architecture"] = arch
    mapping["uart"] = uart
    for name, fn in [
        ("memory_map", gen_memory_map_png),
        ("sim_cycle", gen_simulation_cycle_png),
        ("modules", gen_module_graph_png),
        ("gateway", gen_gateway_overview_png),
    ]:
        p = ASSETS / f"{name}.png"
        fn(p)
        mapping[name] = p
    return mapping


@dataclass
class Section:
    level: int
    title: str
    lines: list[str]
    source: str


def parse_markdown_file(path: Path, source_label: str) -> list[Section]:
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = text.splitlines()
    sections: list[Section] = []
    preamble: list[str] = []
    cur_level = 0
    cur_title = ""
    cur_lines: list[str] = []

    def flush():
        nonlocal cur_lines, cur_title, cur_level
        if cur_title or cur_lines:
            sections.append(Section(cur_level, cur_title or "(без заголовка)", list(cur_lines), source_label))
        cur_lines = []

    for line in lines:
        if line.strip() == "---":
            continue
        m = re.match(r"^(#{1,6})\s+(.*)$", line)
        if m:
            flush()
            cur_level = len(m.group(1))
            cur_title = m.group(2).strip()
            continue
        if not cur_title and not sections:
            preamble.append(line)
            continue
        cur_lines.append(line)

    flush()
    if preamble:
        sections.insert(0, Section(0, "Преамбула", preamble, source_label))
    return sections


def lines_to_bullets(lines: list[str]) -> list[str]:
    bullets: list[str] = []
    in_table = False
    table_buf: list[str] = []

    def flush_table():
        nonlocal table_buf, in_table
        if not table_buf:
            in_table = False
            return
        for row in table_buf:
            row = row.strip()
            if re.match(r"^\|?[\s\-:|]+\|?$", row):
                continue
            cells = [c.strip() for c in row.split("|") if c.strip()]
            if len(cells) >= 2:
                bullets.append(" — ".join(cells[:6]))
            elif cells:
                bullets.append(cells[0])
        table_buf = []
        in_table = False

    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if stripped.startswith("```"):
            flush_table()
            i += 1
            code: list[str] = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code.append(lines[i])
                i += 1
            if i < len(lines):
                i += 1
            if code:
                bullets.append("```")
                bullets.extend(code)
                bullets.append("```")
            continue
        if "|" in stripped and stripped.startswith("|"):
            if not in_table:
                flush_table()
                in_table = True
            table_buf.append(stripped)
            i += 1
            continue
        flush_table()
        if not stripped:
            i += 1
            continue
        if re.match(r"^[\*\-•]\s+", stripped) or re.match(r"^\d+\.\s+", stripped):
            bullets.append(re.sub(r"^[\*\-•]\s+|^\d+\.\s+", "", stripped))
        else:
            wrapped = textwrap.fill(stripped, 100, break_long_words=True)
            for para in wrapped.split("\n"):
                if para.strip():
                    bullets.append(para.strip())
        i += 1
    flush_table()
    return bullets


def normalize_bullets(raw: list[str]) -> list[str]:
    res: list[str] = []
    for b in raw:
        if len(b) > MAX_BULLET_CHARS:
            res.extend(textwrap.wrap(b, MAX_BULLET_CHARS, break_long_words=True))
        else:
            res.append(b)
    return res


def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = add_blank_slide(prs)
    set_bg(slide, C_COVER)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(9), Inches(1.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    st = tf.add_paragraph()
    st.text = subtitle
    st.font.size = Pt(16)
    st.font.color.rgb = RGBColor(0xBB, 0xD4, 0xF0)
    st.space_before = Pt(10)


def add_picture_slide(prs: Presentation, title: str, png: Path, foot: str):
    slide = add_blank_slide(prs)
    set_bg(slide, C_BG)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.35), Inches(9.1), Inches(0.65))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = C_TITLE
    if png.is_file():
        slide.shapes.add_picture(str(png), Inches(0.35), Inches(1.05), width=Inches(9.3))
    fb = slide.shapes.add_textbox(Inches(0.45), Inches(6.95), Inches(9.1), Inches(0.45))
    fp = fb.text_frame.paragraphs[0]
    fp.text = foot
    fp.font.size = Pt(10)
    fp.font.color.rgb = C_MUTED


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], foot: str):
    slide = add_blank_slide(prs)
    set_bg(slide, C_BG)
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.35), Inches(9.1), Inches(0.75))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title[:200]
    tp.font.size = Pt(20)
    tp.font.bold = True
    tp.font.color.rgb = C_TITLE

    body = slide.shapes.add_textbox(Inches(0.45), Inches(1.1), Inches(9.1), Inches(5.75))
    bf = body.text_frame
    bf.word_wrap = True
    for j, line in enumerate(bullets):
        para = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
        para.text = line
        para.font.size = Pt(11) if len(line) < 90 else Pt(10)
        para.font.color.rgb = C_TEXT
        para.space_after = Pt(3)
        para.level = 0
    fb = slide.shapes.add_textbox(Inches(0.45), Inches(6.95), Inches(9.1), Inches(0.45))
    fp = fb.text_frame.paragraphs[0]
    fp.text = foot
    fp.font.size = Pt(9)
    fp.font.color.rgb = C_MUTED


def add_code_slide(prs: Presentation, title: str, code_lines: list[str], foot: str):
    slide = add_blank_slide(prs)
    set_bg(slide, RGBColor(0xFA, 0xFA, 0xFA))
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.35), Inches(9.1), Inches(0.55))
    tb.text_frame.paragraphs[0].text = title[:180]
    tb.text_frame.paragraphs[0].font.size = Pt(16)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = C_TITLE

    chunk = code_lines[:MAX_CODE_LINES]
    text = "\n".join(chunk)
    if len(code_lines) > MAX_CODE_LINES:
        text += f"\n… (+{len(code_lines) - MAX_CODE_LINES} строк)"

    body = slide.shapes.add_textbox(Inches(0.35), Inches(0.95), Inches(9.3), Inches(6.2))
    bf = body.text_frame
    bf.word_wrap = True
    p = bf.paragraphs[0]
    p.text = text
    p.font.name = "Courier New"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x1B, 0x1B, 0x1B)

    fb = slide.shapes.add_textbox(Inches(0.45), Inches(6.95), Inches(9.1), Inches(0.45))
    fb.text_frame.paragraphs[0].text = foot
    fb.text_frame.paragraphs[0].font.size = Pt(9)
    fb.text_frame.paragraphs[0].font.color.rgb = C_MUTED


def emit_section(prs: Presentation, sec: Section, images: dict[str, Path]):
    foot = f"Источник: {sec.source}"
    raw_bullets = lines_to_bullets(sec.lines)
    parts: list[tuple[str, list[str]]] = []
    buf: list[str] = []
    in_code = False
    code_acc: list[str] = []

    for b in raw_bullets:
        if b == "```":
            if in_code:
                if code_acc:
                    parts.append(("code", code_acc))
                    code_acc = []
                in_code = False
            else:
                if buf:
                    parts.append(("text", buf))
                    buf = []
                in_code = True
            continue
        if in_code:
            code_acc.append(b)
        else:
            buf.append(b)
    if buf:
        parts.append(("text", buf))
    if code_acc:
        parts.append(("code", code_acc))

    title_base = sec.title

    for kind, payload in parts:
        if kind == "text":
            bullets = normalize_bullets(payload)
            for i in range(0, len(bullets), MAX_BULLETS_PER_SLIDE):
                chunk = bullets[i : i + MAX_BULLETS_PER_SLIDE]
                suffix = " (продолж.)" if i else ""
                add_bullet_slide(prs, f"{title_base}{suffix}", chunk, foot)
        else:
            for i in range(0, len(payload), MAX_CODE_LINES):
                add_code_slide(
                    prs,
                    f"{title_base} — код",
                    payload[i : i + MAX_CODE_LINES],
                    foot,
                )

    tl = sec.title.lower()
    contextual: list[tuple[str, Path, str]] = []
    if re.search(r"6\.\s+архитектурные|архитектурные требования", tl):
        contextual.append(("Архитектура симулятора (SVG)", images["architecture"], "docs/architecture.svg"))
    if "структура проекта" in tl:
        contextual.append(("Каталог модулей src/ (схема)", images["modules"], "Сгенерировано"))
    if "карта памяти" in tl:
        contextual.append(("Карта памяти", images["memory_map"], "По IMPLEMENTATION.md"))
    if "цикл симуляции" in tl or "порядок вызовов" in tl:
        contextual.append(("Цикл симуляции", images["sim_cycle"], "IMPLEMENTATION.md"))
    if re.search(r"4\.5", tl) and "uart" in tl and "gateway" not in sec.source.lower():
        contextual.append(("UART — диаграмма", images["uart"], "uart_diagram.svg"))
    if "архитектура с поддержкой gdb" in tl or (tl.startswith("2)") and "архитектур" in tl):
        contextual.append(("Онлайн-лаборатория — обзор", images["gateway"], "docs/gateway.md"))

    for pic_title, pic_path, pic_foot in contextual:
        if pic_path.is_file():
            add_picture_slide(prs, pic_title, pic_path, f"{foot} · {pic_foot}")


def build_deck():
    images = generate_all_images()
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "STM32 Emulator — материалы проекта",
        "Полный текст Markdown + рисунки (SVG→PNG и сгенерированные схемы)",
    )

    add_bullet_slide(
        prs,
        "Рисунки в презентации",
        [
            "Архитектура — docs/architecture.svg (после раздела про архитектурные требования в ТЗ).",
            "Модули src/ — блок-схема (после «Структура проекта» в IMPLEMENTATION).",
            "Карта памяти — PNG-таблица (после одноимённого раздела).",
            "Цикл симуляции — после раздела «Цикл симуляции / порядок вызовов».",
            "UART — uart_diagram.svg (после § 4.5 UART в README).",
            "Онлайн-лаборатория — схема по gateway.md (после раздела об архитектуре с GDB).",
        ],
        "Скрипт: presentation/build_full_md_deck.py",
    )

    for path, label in MD_SOURCES:
        if not path.is_file():
            continue
        add_bullet_slide(
            prs,
            f"Документ: {label}",
            [f"Файл: {path.relative_to(ROOT)}", "Далее — разделы по заголовкам # …"],
            "",
        )
        for sec in parse_markdown_file(path, label):
            if not sec.lines and not sec.title:
                continue
            emit_section(prs, sec, images)

    add_title_slide(prs, "Конец", "Вопросы?")

    prs.save(str(OUT_PPTX))
    print("Saved:", OUT_PPTX)


if __name__ == "__main__":
    build_deck()
