#!/usr/bin/env python3
"""
Презентация для защиты проекта STM32F103 Simulator.

Дизайн (по агент-скиллу pptx):
- Тёмная палитра «embedded / oscilloscope»: navy + amber + mint.
- Визуальный мотив: метка адреса (моноширинно) вверху каждого слайда
  и тонкая амбер-полоса слева — повторяется на ВСЕХ контентных слайдах.
- На каждом слайде — визуальный элемент: фигуры, большие цифры, иконки,
  картинки или код-панель. Без скучных «заголовок + буллеты».

Запуск:
  python3 presentation/build_pro_deck.py
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
HERE = Path(__file__).resolve().parent
OUT = HERE / "Защита_проект_STM32.pptx"

# --- ПАЛИТРА (embedded / oscilloscope) ---------------------------------------
NAVY        = RGBColor(0x0B, 0x1B, 0x3B)   # фон тёмных слайдов (доминанта)
NAVY_DEEP   = RGBColor(0x07, 0x12, 0x28)
PANEL       = RGBColor(0xF3, 0xF5, 0xF9)   # фон контентных слайдов
PANEL_SOFT  = RGBColor(0xE8, 0xEC, 0xF3)
INK         = RGBColor(0x0E, 0x1B, 0x32)   # тёмный текст на светлом
INK_SOFT    = RGBColor(0x4A, 0x55, 0x6B)
AMBER       = RGBColor(0xF2, 0xA6, 0x5A)   # акцент 1 (LED-янтарь)
AMBER_SOFT  = RGBColor(0xFF, 0xCB, 0x91)
MINT        = RGBColor(0x4E, 0xC9, 0xB0)   # акцент 2 (терминал)
ICEBLUE     = RGBColor(0xC2, 0xD4, 0xEE)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DIM_LIGHT   = RGBColor(0x8A, 0x95, 0xA8)
TERMINAL_BG = RGBColor(0x10, 0x1A, 0x2E)

# --- РАЗМЕРЫ -----------------------------------------------------------------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- УТИЛИТЫ -----------------------------------------------------------------
def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, rgb):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb


def _rect(slide, x, y, w, h, rgb, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    if not line:
        s.line.fill.background()
    return s


def _txt(
    slide, x, y, w, h, text, *,
    font="Calibri", size=14, bold=False, italic=False, color=INK,
    align=PP_ALIGN.LEFT, line_spacing=1.15,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box, tf, p


def _bullets(slide, x, y, w, h, items, *,
             color=INK, size=15, bullet_color=AMBER, font="Calibri"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.line_spacing = 1.18
        para.space_after = Pt(4)
        bullet_run = para.add_run()
        bullet_run.text = "■  "
        bullet_run.font.name = "Consolas"
        bullet_run.font.size = Pt(size - 2)
        bullet_run.font.color.rgb = bullet_color
        text_run = para.add_run()
        text_run.text = it
        text_run.font.name = font
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = color
    return box


def _motif(slide, addr_tag: str, page_no: int, total: int, on_dark=False):
    """Визуальный мотив: левая амбер-полоса + метка адреса вверху + футер."""
    _rect(slide, Inches(0), Inches(0), Inches(0.10), SLIDE_H, AMBER)
    if on_dark:
        tag_color = AMBER_SOFT
        foot_color = DIM_LIGHT
    else:
        tag_color = INK_SOFT
        foot_color = INK_SOFT
    _txt(slide, Inches(10.4), Inches(0.30), Inches(2.7), Inches(0.35),
         addr_tag, font="Consolas", size=12, color=tag_color, align=PP_ALIGN.RIGHT)
    _txt(slide, Inches(0.4), Inches(7.10), Inches(8), Inches(0.3),
         "STM32F103 Simulator · ИТМО · ПрВС 2026",
         font="Consolas", size=10, color=foot_color)
    _txt(slide, Inches(11.6), Inches(7.10), Inches(1.5), Inches(0.3),
         f"{page_no:02d} / {total:02d}",
         font="Consolas", size=10, color=foot_color, align=PP_ALIGN.RIGHT)


# --- ЛЭЙАУТЫ -----------------------------------------------------------------
def slide_cover(prs, title: str, subtitle: str, tagline: str):
    s = _blank(prs)
    _bg(s, NAVY)
    _rect(s, Inches(0), Inches(0), Inches(0.10), SLIDE_H, AMBER)
    _rect(s, Inches(0.4), Inches(2.0), Inches(0.18), Inches(2.5), MINT)
    _txt(s, Inches(0.85), Inches(1.2), Inches(11), Inches(0.6),
         "// embedded systems · 2026",
         font="Consolas", size=14, color=AMBER_SOFT)
    _txt(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(2.0),
         title, font="Cambria", size=54, bold=True, color=WHITE,
         line_spacing=1.05)
    _txt(s, Inches(0.85), Inches(4.2), Inches(11.6), Inches(0.7),
         subtitle, font="Calibri", size=22, color=ICEBLUE)
    _rect(s, Inches(0.85), Inches(5.4), Inches(0.6), Inches(0.04), AMBER)
    _txt(s, Inches(0.85), Inches(5.55), Inches(11.6), Inches(0.5),
         tagline, font="Consolas", size=13, color=DIM_LIGHT)
    return s


def slide_section(prs, num: str, kicker: str, title: str, page: int, total: int):
    s = _blank(prs)
    _bg(s, NAVY)
    _rect(s, Inches(0), Inches(0), Inches(0.10), SLIDE_H, AMBER)
    _txt(s, Inches(0.85), Inches(2.6), Inches(3), Inches(2.5),
         num, font="Cambria", size=180, bold=True, color=AMBER,
         line_spacing=1.0)
    _txt(s, Inches(4.6), Inches(3.0), Inches(8), Inches(0.5),
         kicker, font="Consolas", size=14, color=MINT)
    _txt(s, Inches(4.6), Inches(3.4), Inches(8), Inches(2),
         title, font="Cambria", size=44, bold=True, color=WHITE,
         line_spacing=1.05)
    _rect(s, Inches(4.6), Inches(5.6), Inches(1.5), Inches(0.04), AMBER)
    _txt(s, Inches(11.6), Inches(7.10), Inches(1.5), Inches(0.3),
         f"{page:02d} / {total:02d}",
         font="Consolas", size=10, color=DIM_LIGHT, align=PP_ALIGN.RIGHT)
    return s


def slide_basic(prs, title, lead, bullets, addr_tag, page, total):
    s = _blank(prs)
    _bg(s, PANEL)
    _motif(s, addr_tag, page, total)
    _txt(s, Inches(0.55), Inches(0.55), Inches(11), Inches(0.7),
         title, font="Cambria", size=32, bold=True, color=INK)
    _rect(s, Inches(0.55), Inches(1.30), Inches(0.7), Inches(0.05), AMBER)
    _txt(s, Inches(0.55), Inches(1.50), Inches(12.2), Inches(1.3),
         lead, font="Calibri", size=17, color=INK, line_spacing=1.30)
    return s


def slide_text(prs, title, lead, bullets, addr_tag, page, total):
    s = slide_basic(prs, title, lead, bullets, addr_tag, page, total)
    _bullets(s, Inches(0.55), Inches(3.0), Inches(12.2), Inches(4.0),
             bullets, color=INK, size=16)
    return s


def slide_stat(prs, title, lead, big_value, big_label,
               sub_bullets, addr_tag, page, total):
    s = slide_basic(prs, title, lead, [], addr_tag, page, total)
    # Картина-калаут: огромная цифра слева, буллеты справа
    _rect(s, Inches(0.55), Inches(3.1), Inches(5.0), Inches(3.6), NAVY)
    _txt(s, Inches(0.7), Inches(3.2), Inches(4.7), Inches(0.5),
         "// stat", font="Consolas", size=12, color=AMBER_SOFT)
    _txt(s, Inches(0.7), Inches(3.55), Inches(4.7), Inches(2.5),
         big_value, font="Cambria", size=110, bold=True, color=AMBER,
         line_spacing=1.0)
    _txt(s, Inches(0.7), Inches(6.05), Inches(4.7), Inches(0.6),
         big_label, font="Calibri", size=18, color=ICEBLUE)
    _bullets(s, Inches(6.0), Inches(3.1), Inches(7), Inches(4.0),
             sub_bullets, color=INK, size=16)
    return s


def slide_two_col(prs, title, lead, left_h, left_lines, right_h, right_lines,
                  addr_tag, page, total):
    s = slide_basic(prs, title, lead, [], addr_tag, page, total)
    # Левая колонка
    _rect(s, Inches(0.55), Inches(3.0), Inches(0.04), Inches(3.8), MINT)
    _txt(s, Inches(0.75), Inches(3.0), Inches(5.5), Inches(0.5),
         left_h, font="Cambria", size=20, bold=True, color=INK)
    _bullets(s, Inches(0.75), Inches(3.5), Inches(5.5), Inches(3.4),
             left_lines, color=INK, size=14, bullet_color=MINT)
    # Правая колонка
    _rect(s, Inches(6.85), Inches(3.0), Inches(0.04), Inches(3.8), AMBER)
    _txt(s, Inches(7.05), Inches(3.0), Inches(5.7), Inches(0.5),
         right_h, font="Cambria", size=20, bold=True, color=INK)
    _bullets(s, Inches(7.05), Inches(3.5), Inches(5.7), Inches(3.4),
             right_lines, color=INK, size=14, bullet_color=AMBER)
    return s


def slide_image(prs, title, lead, png: Path, caption,
                addr_tag, page, total, *, height=Inches(3.7)):
    s = slide_basic(prs, title, lead, [], addr_tag, page, total)
    if png.is_file():
        # Расположение: по центру под текстом
        # Высота фиксирована, ширина рассчитывается автоматом
        s.shapes.add_picture(str(png), Inches(0.7), Inches(2.95), height=height)
    if caption:
        _txt(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.6),
             caption, font="Calibri", size=14, italic=True, color=INK_SOFT)
    return s


def slide_code(prs, title, explain, lines, addr_tag, page, total):
    s = _blank(prs)
    _bg(s, PANEL)
    _motif(s, addr_tag, page, total)
    _txt(s, Inches(0.55), Inches(0.55), Inches(11), Inches(0.7),
         title, font="Cambria", size=30, bold=True, color=INK)
    _rect(s, Inches(0.55), Inches(1.28), Inches(0.7), Inches(0.05), AMBER)
    _txt(s, Inches(0.55), Inches(1.50), Inches(12.2), Inches(0.9),
         explain, font="Calibri", size=15, color=INK, line_spacing=1.25)
    _rect(s, Inches(0.55), Inches(2.65), Inches(12.2), Inches(4.2),
          TERMINAL_BG)
    code_text = "\n".join(lines)
    _txt(s, Inches(0.85), Inches(2.85), Inches(11.7), Inches(3.9),
         code_text, font="Consolas", size=14, color=MINT, line_spacing=1.30)
    _txt(s, Inches(0.85), Inches(2.65), Inches(11.7), Inches(0.3),
         "$ shell", font="Consolas", size=10, color=AMBER_SOFT)
    return s


def slide_table(prs, title, lead, headers, rows,
                addr_tag, page, total, col_widths=None):
    s = slide_basic(prs, title, lead, [], addr_tag, page, total)
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [12.2 / n_cols] * n_cols
    x = 0.55
    y = 3.0
    row_h = 0.45
    # Заголовок
    cur_x = x
    for i, h in enumerate(headers):
        w = col_widths[i]
        _rect(s, Inches(cur_x), Inches(y), Inches(w), Inches(row_h), NAVY)
        _txt(s, Inches(cur_x + 0.12), Inches(y + 0.08),
             Inches(w - 0.2), Inches(row_h),
             h, font="Cambria", size=14, bold=True, color=WHITE)
        cur_x += w
    y += row_h
    # Строки
    for ri, row in enumerate(rows):
        cur_x = x
        bg = PANEL_SOFT if ri % 2 == 0 else PANEL
        _rect(s, Inches(x), Inches(y), Inches(sum(col_widths)),
              Inches(row_h), bg)
        for ci, cell in enumerate(row):
            w = col_widths[ci]
            font = "Consolas" if ci == 0 else "Calibri"
            color = AMBER if ci == 0 else INK
            size = 13 if ci == 0 else 13
            _txt(s, Inches(cur_x + 0.12), Inches(y + 0.08),
                 Inches(w - 0.2), Inches(row_h),
                 cell, font=font, size=size, color=color)
            cur_x += w
        y += row_h
    return s


def slide_iconrow(prs, title, lead, items, addr_tag, page, total):
    """3 «карточки» с цветным кругом (иконка-инициал) + заголовок + текст."""
    s = slide_basic(prs, title, lead, [], addr_tag, page, total)
    n = len(items)
    total_w = 12.2
    gap = 0.4
    card_w = (total_w - gap * (n - 1)) / n
    y = 3.0
    card_h = 3.7
    accents = [AMBER, MINT, ICEBLUE, AMBER_SOFT]
    for i, (initial, head, body) in enumerate(items):
        x = 0.55 + i * (card_w + gap)
        _rect(s, Inches(x), Inches(y), Inches(card_w), Inches(card_h), PANEL_SOFT)
        _rect(s, Inches(x), Inches(y), Inches(0.05), Inches(card_h),
              accents[i % len(accents)])
        # Круг с инициалом
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(x + 0.35), Inches(y + 0.35),
                                  Inches(0.95), Inches(0.95))
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY
        circ.line.fill.background()
        _txt(s, Inches(x + 0.35), Inches(y + 0.40),
             Inches(0.95), Inches(0.95),
             initial, font="Cambria", size=28, bold=True,
             color=accents[i % len(accents)], align=PP_ALIGN.CENTER)
        _txt(s, Inches(x + 0.30), Inches(y + 1.45),
             Inches(card_w - 0.5), Inches(0.6),
             head, font="Cambria", size=18, bold=True, color=INK)
        _txt(s, Inches(x + 0.30), Inches(y + 2.05),
             Inches(card_w - 0.5), Inches(card_h - 2.2),
             body, font="Calibri", size=13, color=INK_SOFT,
             line_spacing=1.30)
    return s


# --- ЗАГРУЗКА КАРТИНОК -------------------------------------------------------
def load_image_generator():
    path = HERE / "build_full_md_deck.py"
    name = "_stm32_deck_images"
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    assert spec.loader
    spec.loader.exec_module(mod)
    return mod.generate_all_images


# --- СБОРКА ------------------------------------------------------------------
def build():
    images = load_image_generator()()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Считаем общее количество слайдов заранее (просто фиксированное число)
    TOTAL = 40
    p = 0

    def nxt():
        nonlocal p
        p += 1
        return p

    # 01 — Cover
    slide_cover(
        prs,
        "Симулятор STM32F103",
        "Курсовой проект · кафедра ВТ · ИТМО",
        "C11 · ARM Cortex-M3 · GDB Remote Serial Protocol · 11 наборов тестов",
    )

    # 02 — Что это
    slide_iconrow(
        prs,
        "Что это в одном кадре",
        "Программа на компьютере, которая ведёт себя как микроконтроллер STM32F103. "
        "Запускает прошивку, исполняет инструкции ARM Cortex-M3, моделирует таймер, UART и прерывания.",
        [
            ("S", "Симулятор",
             "Чистый C11. Один поток. Детерминированное поведение: "
             "одинаковый вход — одинаковый результат."),
            ("D", "Отладка",
             "Встроенный GDB-сервер по протоколу RSP, "
             "порт 3333. Подключаемся arm-none-eabi-gdb."),
            ("T", "Тесты",
             "11 наборов модульных и интеграционных тестов. "
             "Запуск: cd src && make test."),
        ],
        "0x08000000  // entry", nxt(), TOTAL,
    )

    # 03 — Зачем
    slide_text(
        prs,
        "Зачем это нужно",
        "Лабораторная работа на железе требует отдельной платы и времени. "
        "Симулятор даёт повторяемый сценарий с пошаговой отладкой и логом — без риска «сжечь» Blue Pill.",
        [
            "Учёба: посмотреть, как взаимодействуют ядро, шина и периферия.",
            "Демо на ноутбуке без подключения USB-программатора.",
            "Воспроизводимая среда: одинаковая на macOS, Linux, Windows-WSL.",
            "Отладка прошивки в attached-режиме через GDB как на железе.",
            "Возможность подавать байты в UART прямо из CLI и смотреть реакцию.",
        ],
        "// motivation", nxt(), TOTAL,
    )

    # 04 — STM32F103 на пальцах
    slide_table(
        prs,
        "STM32F103C8T6 vs наш симулятор",
        "Реальная плата Blue Pill и наш программный двойник.",
        ["Параметр", "Реальный чип", "Наш симулятор"],
        [
            ["Ядро",        "Cortex-M3 @ 72 МГц",       "Cortex-M3 (Thumb-2, без частоты)"],
            ["Flash",       "64 КБ @ 0x08000000",        "64 КБ @ 0x08000000"],
            ["SRAM",        "20 КБ @ 0x20000000",        "20 КБ @ 0x20000000"],
            ["Прерывания",  "до 60+ источников",          "43 IRQ + SysTick"],
            ["TIM2",        "регистры CR1/PSC/ARR/...",   "CR1/DIER/SR/CNT/PSC/ARR"],
            ["USART1",      "TX/RX/IRQ",                  "TX/RX/IRQ (упрощённо)"],
            ["GPIO",        "5 портов A–E",               "5 портов A–E"],
            ["Отладка",     "SWD/JTAG",                  "GDB RSP по TCP, порт 3333"],
        ],
        "0x40000000  // mmio", nxt(), TOTAL,
        col_widths=[3.0, 4.6, 4.6],
    )

    # 05 — ТЗ кратко (большая цифра)
    slide_stat(
        prs,
        "Что просило ТЗ",
        "Из README.md: групповой проект до 6 человек. Симулятор должен иметь ядро, память, NVIC, "
        "таймер, UART и систему тестирования.",
        "6", "человек в команде",
        [
            "Архитектор / разработчик ядра CPU.",
            "Память и загрузка прошивки.",
            "NVIC и обработка прерываний.",
            "Таймер.",
            "UART и взаимодействие с консолью.",
            "Инженер по тестированию и сборке.",
        ],
        "// README §8", nxt(), TOTAL,
    )

    # 06 — Стек технологий
    slide_two_col(
        prs,
        "Технологический стек",
        "Намеренно простой: чистый C, минимум зависимостей, всё собирается одной командой make.",
        "Симулятор",
        [
            "Язык: C11, gcc/clang.",
            "Сборка: GNU Make.",
            "Зависимости: только стандартная C-библиотека.",
            "TCP-сервер для GDB через POSIX sockets.",
        ],
        "Прошивка-пример",
        [
            "Кросс-компилятор arm-none-eabi-gcc.",
            "Свой linker script (link.ld) под адреса 0x08000000 / 0x20000000.",
            "Стартап на C: Reset_Handler + vector_table.",
            "objcopy → firmware.bin для загрузки в симулятор.",
        ],
        "// toolchain", nxt(), TOTAL,
    )

    # ----- РАЗДЕЛ 1: АРХИТЕКТУРА ----------------------------------
    slide_section(prs, "01", "// section",
                  "Архитектура симулятора", nxt(), TOTAL)

    # 08 — общая схема (картинка)
    slide_image(
        prs,
        "Общая схема (architecture.svg)",
        "Ядро исполняет инструкции и обращается к шине. Шина по адресу решает, "
        "куда направить запрос: Flash, SRAM или периферия. NVIC принимает запросы прерываний.",
        images["architecture"],
        "Стрелки — направление данных. Все обращения от ядра — только через шину; "
        "периферия не знает о ядре напрямую.",
        "// docs/architecture.svg", nxt(), TOTAL,
        height=Inches(3.7),
    )

    # 09 — модули src/
    slide_image(
        prs,
        "Модули каталога src/",
        "Каждая папка — изолированный модуль с публичным заголовком (.h) и реализацией (.c). "
        "Нижний слой — периферия и отладка, средний — шина и память, вверху — ядро.",
        images["modules"],
        "Связи направлены сверху вниз: симулятор оркестрирует, ядро говорит с шиной, шина — с устройствами.",
        "// src/", nxt(), TOTAL,
        height=Inches(3.5),
    )

    # 10 — карта памяти
    slide_image(
        prs,
        "Карта памяти",
        "Те же адреса, что в datasheet STM32F103: прошивка обращается к регистрам так же, как на железе.",
        images["memory_map"],
        "Адрес 0x00000000 — алиас Flash для чтения вектор-таблицы. "
        "Регистры периферии живут с 0x40000000.",
        "0x08000000 / 0x20000000", nxt(), TOTAL,
        height=Inches(3.2),
    )

    # 11 — Reset
    slide_stat(
        prs,
        "Что происходит после reset",
        "В Cortex-M3 первые два слова Flash — это начальное значение SP и адрес Reset_Handler. "
        "Симулятор в core_reset загружает их через шину.",
        "0x00",
        "первое слово → SP, второе слово (0x04) → PC",
        [
            "SP ← bus_read(0x00000000, 4).",
            "PC ← bus_read(0x00000004, 4) & ~1 (сбрасываем Thumb-бит).",
            "Сохраняются указатели bus и nvic, остальное обнуляется.",
            "После — full reset: NVIC, память, очередь событий, профайлер.",
            "Если PC указывает не во Flash — будет bus fault (это диагностика).",
        ],
        "core_reset()", nxt(), TOTAL,
    )

    # 12 — цикл шага
    slide_image(
        prs,
        "Цикл одного шага симуляции",
        "Порядок зафиксирован в simulator_step. Перестановка ломает детерминизм, поэтому изменения там запрещены.",
        images["sim_cycle"],
        "1) cycle++ → 2) tick всех периферий → 3) dispatch событий → 4) core_step → 5) check breakpoints.",
        "simulator_step()", nxt(), TOTAL,
        height=Inches(2.8),
    )

    # ----- РАЗДЕЛ 2: ЯДРО ----------------------------------------
    slide_section(prs, "02", "// section",
                  "Ядро Cortex-M3", nxt(), TOTAL)

    # 14 — Состояние ядра
    slide_two_col(
        prs,
        "Состояние ядра (CoreState)",
        "Структура из core.h хранит весь регистровый файл и служебные флаги. На неё опираются декодер, "
        "вход/выход прерываний и отладчик через GDB.",
        "Регистры",
        [
            "r[16] — общие регистры R0–R15 (включая SP, LR, PC).",
            "xpsr — флаги N/Z/C/V в битах 31–28.",
            "thumb_mode — всегда true для Cortex-M3.",
        ],
        "Служебные поля",
        [
            "interruptible — разрешены ли прерывания.",
            "current_irq — кодируется IRQ либо core-исключение (0x8000 | n).",
            "cycles — счётчик тактов для метрик и тестов.",
        ],
        "core/core.h", nxt(), TOTAL,
    )

    # 15 — декодер
    slide_stat(
        prs,
        "Декодер инструкций Thumb",
        "В core.c таблица instr_table сопоставляет 16-битные инструкции по mask/pattern. "
        "Реализованы все ключевые «форматы» из ARM Architecture Reference Manual.",
        "16",
        "форматов Thumb-1 + 32-битные Thumb-2",
        [
            "Сдвиги (1), ADD/SUB (2), immediate-8 (3).",
            "ALU register (4), Hi/BX (5), PC-relative LDR (6).",
            "Reg-offset и imm-offset load/store (7–10), SP-relative (11).",
            "ADR/SP±, PUSH/POP (12–14), BCC, SVC, B (16–18).",
            "Условные ветви через condition_passed (флаги xPSR).",
        ],
        "instr_table[]", nxt(), TOTAL,
    )

    # 16 — Thumb-2
    slide_text(
        prs,
        "Расширение Thumb-2 (32-битные)",
        "Если первые биты инструкции попадают в область 32-битных Thumb-2, выполнение уходит в execute_32bit. "
        "Это критично для прошивок, собранных современным arm-none-eabi-gcc.",
        [
            "BL — 32-битный branch with link для вызовов функций.",
            "MOVW / MOVT — загрузка 16-битных immediate в регистр.",
            "modified-immediate dataproc — арифметика и логика с константами.",
            "wide LDR/STR — обращения с большим смещением.",
            "знаковые загрузки и UMULL/SMULL для 64-битных операций.",
        ],
        "execute_32bit()", nxt(), TOTAL,
    )

    # 17 — EXC_RETURN
    slide_table(
        prs,
        "Вход и выход из прерывания",
        "Cortex-M3 кладёт в стек кадр и подменяет LR специальным магическим значением EXC_RETURN.",
        ["Значение LR", "Что означает", "Куда вернуться"],
        [
            ["0xFFFFFFF1", "возврат в Handler-mode", "MSP, привилегированно"],
            ["0xFFFFFFF9", "возврат в Thread-mode + MSP", "main-стек"],
            ["0xFFFFFFFD", "возврат в Thread-mode + PSP", "process-стек"],
            ["IS_EXC_RETURN", "(LR & 0xFFFFFFF0) == 0xFFFFFFF0", "детектируется в BX/POP"],
        ],
        "EXC_RETURN", nxt(), TOTAL,
        col_widths=[3.5, 5.4, 3.3],
    )

    # ----- РАЗДЕЛ 3: NVIC ----------------------------------------
    slide_section(prs, "03", "// section",
                  "NVIC — контроллер прерываний", nxt(), TOTAL)

    # 19 — NVIC stats
    slide_stat(
        prs,
        "NVIC: что внутри",
        "Структура NVIC хранит массивы pending/active/enabled/priority и текущий приоритет. "
        "Логика выбора IRQ — функция nvic_get_pending_irq.",
        "43",
        "линий внешних прерываний + SysTick",
        [
            "Приоритет: меньшее число = важнее.",
            "Выбирается enabled+pending IRQ с наименьшим priority < current_priority.",
            "nvic_acknowledge — clear pending, set active.",
            "nvic_complete — clear active и пересчёт current_priority.",
            "SysTick живёт отдельным набором полей, проверяется в core_step.",
        ],
        "nvic.h", nxt(), TOTAL,
    )

    # 20 — NVIC MMIO
    slide_table(
        prs,
        "Регистры NVIC через шину",
        "В nvic_bus.c регион 0xE000E100 + 0x400 — стандартный MMIO интерфейс Cortex-M к NVIC.",
        ["Адрес (от базы)", "Регистр", "Назначение"],
        [
            ["0x000",  "ISER0/1",        "set-enable IRQ (write-1)"],
            ["0x080",  "ICER0/1",        "clear-enable IRQ"],
            ["0x100",  "ISPR0/1",        "set-pending IRQ"],
            ["0x180",  "ICPR0/1",        "clear-pending IRQ"],
            ["0x200",  "IABR0/1",        "active bits (только чтение)"],
            ["0x300+", "IPR0..n",        "приоритеты (по 8 бит на IRQ)"],
        ],
        "0xE000E100", nxt(), TOTAL,
        col_widths=[2.8, 2.6, 6.8],
    )

    # ----- РАЗДЕЛ 4: ПАМЯТЬ ----------------------------------
    slide_section(prs, "04", "// section",
                  "Память и шина", nxt(), TOTAL)

    # 22 — Memory
    slide_two_col(
        prs,
        "Memory: Flash и SRAM",
        "memory.c использует два байтовых массива и пары read_le/write_le для little-endian доступа.",
        "Flash · 64 КБ",
        [
            "База 0x08000000.",
            "Только чтение: попытка записи → STATUS_ERROR.",
            "memory_load_binary читает .bin от 0 до 64 КБ.",
            "На алиас 0x00000000 виден тот же образ.",
        ],
        "SRAM · 20 КБ",
        [
            "База 0x20000000.",
            "Чтение/запись 1, 2 или 4 байта — иное возвращает 0.",
            "OOB-запись → STATUS_INVALID_ADDRESS.",
            "memory_reset очищает SRAM, Flash сохраняется.",
        ],
        "memory.c", nxt(), TOTAL,
    )

    # 23 — Bus
    slide_stat(
        prs,
        "Bus — диспетчер адресов",
        "bus.c хранит до 24 BusRegion (base, size, ctx, read, write). "
        "Линейный поиск находит нужный регион и вычисляет offset = addr − base.",
        "24",
        "максимум регистрируемых регионов",
        [
            "bus_register_region добавляет регион — без авто-проверки пересечений.",
            "Read miss → 0 + сообщение в stderr (диагностика).",
            "Write miss → STATUS_INVALID_ADDRESS + stderr.",
            "Размер обращений: 1, 2 или 4 байта.",
            "Ядро ходит в память ТОЛЬКО через bus_read / bus_write.",
        ],
        "bus.c", nxt(), TOTAL,
    )

    # ----- РАЗДЕЛ 5: ПЕРИФЕРИЯ ----------------------------------
    slide_section(prs, "05", "// section",
                  "Периферия", nxt(), TOTAL)

    # 25 — TIM2
    slide_table(
        prs,
        "TIM2 — таймер общего назначения",
        "В timer.c CNT вычисляется не каждый такт, а через event-driven модель: переполнение планируется в EventQueue.",
        ["Смещение", "Регистр", "Что делает"],
        [
            ["0x00", "CR1",  "включение/режим (бит CEN)"],
            ["0x0C", "DIER", "разрешение прерываний (UIE)"],
            ["0x10", "SR",   "флаги; UIF — обновление; пишется как W0C"],
            ["0x24", "CNT",  "значение счётчика"],
            ["0x28", "PSC",  "прескалер; период = (PSC+1) тактов"],
            ["0x2C", "ARR",  "auto-reload; 0 или CEN=0 → CNT читается как 0"],
        ],
        "0x40000000", nxt(), TOTAL,
        col_widths=[1.7, 1.7, 8.8],
    )

    # 26 — USART1
    slide_two_col(
        prs,
        "USART1 — асинхронный приёмопередатчик",
        "В uart.c регистры на 0x40013800. TX «мгновенный»: запись в DR ставит tx_pending, следующий tick завершает передачу.",
        "Передача (TX)",
        [
            "Запись в DR при UE=1 → сбрасывает TXE/TC.",
            "uart_tick: восстанавливает TXE|TC, вызывает output_fn.",
            "При TXEIE+UE — nvic_set_pending после передачи.",
        ],
        "Приём (RX)",
        [
            "Кольцо на 16 байт.",
            "RXNE сбрасывается, когда буфер пуст.",
            "uart_incoming_char и RXNEIE+UE → IRQ.",
        ],
        "0x40013800", nxt(), TOTAL,
    )

    # 27 — UART diagram
    slide_image(
        prs,
        "UART — реальная диаграмма",
        "На остановке симуляции gdb_stub пишет файл uart_diagram.html. Полезно показать «что реально шёл байт за байтом».",
        images["uart"],
        "Каждый столбик — символ на TX. Ось снизу — такты симуляции.",
        "uart_diagram.svg", nxt(), TOTAL,
        height=Inches(2.6),
    )

    # 28 — GPIO
    slide_table(
        prs,
        "GPIO — пять портов A…E",
        "В gpio.c каждый порт — отдельный экземпляр Peripheral с одинаковыми регистрами.",
        ["Порт", "База", "Что внутри"],
        [
            ["GPIOA", "0x40010800", "IDR, ODR, BSRR, BRR, CRL, CRH, LCKR (stub)"],
            ["GPIOB", "0x40010C00", "то же"],
            ["GPIOC", "0x40011000", "PC13 — встроенный светодиод Blue Pill"],
            ["GPIOD", "0x40011400", "то же"],
            ["GPIOE", "0x40011800", "то же"],
        ],
        "GPIO[A..E]", nxt(), TOTAL,
        col_widths=[1.6, 2.5, 8.1],
    )

    # 29 — RCC + SysTick
    slide_two_col(
        prs,
        "RCC и SysTick",
        "Тактирование и системный таймер реализованы как «честные заглушки»: регистры читаются и пишутся корректно, но без модели частоты.",
        "RCC · 0x40021000",
        [
            "10 регистров: CR, CFGR, CIR, ..., CSR.",
            "SWS отслеживает SW (выбор источника).",
            "READY-биты повторяют ENABLE-биты.",
            "Доступ только по 4-байтному выравниванию.",
        ],
        "SysTick · 0xE000E010",
        [
            "16 байт: CSR, RVR, CVR, CALIB.",
            "24-битная маска значений RVR/CVR.",
            "COUNTFLAG сбрасывается при чтении CSR.",
            "TICKINT → отдельная линия прерывания через NVIC.",
        ],
        "// stub-ish", nxt(), TOTAL,
    )

    # 30 — Peripheral interface
    slide_code(
        prs,
        "Как добавить новую периферию (vtable в чистом C)",
        "Новая периферия — это структура Peripheral с указателями на 4 функции. "
        "Симулятор регистрирует её на шине; существующий код не трогаем (Open/Closed).",
        [
            "typedef struct {",
            "    void*    ctx;                                                              /* ваше состояние */",
            "    uint32_t (*read )(void* ctx, uint32_t offset, uint8_t size);                /* чтение регистра */",
            "    Status   (*write)(void* ctx, uint32_t offset, uint32_t value, uint8_t size);/* запись регистра */",
            "    void     (*tick )(void* ctx);                                              /* такт симуляции */",
            "    void     (*reset)(void* ctx);                                              /* сброс */",
            "} Peripheral;",
            "",
            "/* В simulator.c: */",
            "simulator_add_peripheral(&sim, MY_DEV_BASE, MY_DEV_SIZE, my_dev_as_peripheral(&dev));",
        ],
        "peripherals/peripheral.h", nxt(), TOTAL,
    )

    # ----- РАЗДЕЛ 6: ОТЛАДКА ------------------------------------
    slide_section(prs, "06", "// section",
                  "Отладка и пользование", nxt(), TOTAL)

    # 32 — CLI
    slide_text(
        prs,
        "CLI — 15 интерактивных команд",
        "В ui.c таблица commands[] перечисляет все команды; ui_run читает строку и диспатчит её на обработчик. "
        "Парсинг адресов через strtoul base=0 — поддерживает 0x-нотацию.",
        [
            "help, quit — справка и выход.",
            "load <путь> — загрузить .bin во Flash.",
            "run / stop — продолжить / приостановить симуляцию до breakpoint.",
            "step [N] — выполнить N инструкций.",
            "reset — full reset (NVIC + память + события).",
            "reg / mem <addr> [count] — регистры и hex-дамп памяти.",
            "break <addr> / delete <addr> — точки останова.",
            "uart <char>, diagram, diagram-save, profile — телеметрия и измерения.",
        ],
        "ui_run()", nxt(), TOTAL,
    )

    # 33 — GDB stub
    slide_table(
        prs,
        "GDB Remote Serial Protocol — что поддерживаем",
        "gdb_stub.c — TCP-сервер на порту 3333. После accept крутится rsp_session; команды бьются на пакеты со стандартным CRC.",
        ["Пакет", "Действие", "Реализация"],
        [
            ["?",          "стоп-причина",              "S05 (trap)"],
            ["g / G",      "все регистры",              "из CoreState.r[]"],
            ["m / M",      "память",                    "bus_read / bus_write"],
            ["c / s",      "continue / step",           "simulator_step + Ctrl+C"],
            ["Z / z",      "программный breakpoint",    "debugger_add_breakpoint"],
            ["qSupported", "возможности",               "PacketSize=1000;qXfer:features:read+"],
            ["qRcmd",      "monitor-команды",           "halt, reset, uart-diagram, uart-svg"],
        ],
        "TCP :3333", nxt(), TOTAL,
        col_widths=[2.0, 4.0, 6.5],
    )

    # 34 — How to run
    slide_code(
        prs,
        "Как запустить и как подключиться",
        "Стандартные команды для защиты: соберите симулятор, прогоните тесты, запустите прошивку, "
        "при желании подключите GDB.",
        [
            "# Сборка",
            "$ cd src",
            "$ make",
            "",
            "# Тесты",
            "$ make test",
            "",
            "# Запуск с прошивкой",
            "$ ./stm32sim ../examples/firmware.bin",
            "",
            "# Запуск с GDB-сервером",
            "$ ./stm32sim ../examples/firmware.bin --gdb 3333",
            "",
            "# В другом терминале:",
            "$ arm-none-eabi-gdb ../examples/firmware.elf",
            "(gdb) target remote :3333",
            "(gdb) break main",
            "(gdb) continue",
        ],
        "make · gdb", nxt(), TOTAL,
    )

    # ----- РАЗДЕЛ 7: ТЕСТЫ ------------------------------------
    slide_section(prs, "07", "// section",
                  "Тестирование", nxt(), TOTAL)

    # 36 — Test stat
    slide_stat(
        prs,
        "11 наборов тестов",
        "Драйвер test_main.c вызывает наборы по очереди и считает pass/fail. "
        "Помощники в test.h: ASSERT_EQ, ASSERT_EQ64, RUN_TEST, TEST_SUITE.",
        "11",
        "файлов tests/test_*.c",
        [
            "test_core — Thumb MOV/ADD/CMP, ветви, LDR/STR, PUSH/POP, MUL.",
            "test_nvic — pending/enable/priority/acknowledge/clear.",
            "test_memory, test_bus — границы и роутинг.",
            "test_timer, test_systick — countdown, прерывания.",
            "test_rcc, test_gpio, test_uart — регистры периферии.",
            "test_debugger, test_integration — TIM2→IRQ, breakpoint halt, UART.",
        ],
        "src/tests/", nxt(), TOTAL,
    )

    # 37 — Examples
    slide_two_col(
        prs,
        "Пример прошивки в examples/",
        "Чтобы было что показывать, рядом лежит простая прошивка под arm-none-eabi-gcc.",
        "Что собирается",
        [
            "startup.c — Reset_Handler + vector_table.",
            "main.c — bare-metal с регистрами RCC/GPIO/TIM2/USART1/NVIC.",
            "link.ld — FLASH 64K, SRAM 20K, ENTRY(Reset_Handler).",
            "Output: firmware.elf и firmware.bin.",
        ],
        "Утилиты",
        [
            "gen_firmware.py — синтетическая Thumb-программа для тестов.",
            "gen_stub_elf.py — минимальный ELF32-ARM для отладчиков.",
            "Запуск: ../src/stm32sim firmware.bin.",
        ],
        "examples/", nxt(), TOTAL,
    )

    # 38 — gateway plan
    slide_image(
        prs,
        "Дорожная карта: онлайн-лаборатория (план)",
        "docs/gateway.md описывает будущее: очередь заданий, воркеры, KeyDB, браузерный клиент и удалённый GDB. "
        "Этого кода в текущем репозитории нет — это honest roadmap.",
        images["gateway"],
        "В ответе на вопрос «а что дальше?»: распределённая обвязка вокруг существующего симулятора.",
        "// docs/gateway.md", nxt(), TOTAL,
        height=Inches(2.8),
    )

    # 39 — Ограничения
    slide_text(
        prs,
        "Что осознанно вне MVP",
        "Чтобы не обманывать проверяющего: вот области, где симулятор намеренно упрощён.",
        [
            "Нет полной модели MPU, FPU и кэшей (в Cortex-M3 их и нет).",
            "RCC и SysTick без реальной частоты — модель дискретных тактов.",
            "UART — без BRR-таймингов, передача завершается за один tick.",
            "Алиньмент проверяется только по размеру, не по выравниванию.",
            "Нет CI/Docker — тесты запускаются локально через make test.",
            "Онлайн-сервис из gateway.md ещё не реализован — только дизайн-документ.",
        ],
        "// honest", nxt(), TOTAL,
    )

    # 40 — Итог + Команда + спасибо? Совмещаем
    s = slide_basic(
        prs,
        "Итог и кому что делать дальше",
        "Готовый детерминированный симулятор STM32F103 с тестами, отладкой по GDB и примером прошивки. "
        "Базис для следующих лабораторных и для распределённой онлайн-обёртки.",
        [],
        "// summary", nxt(), TOTAL,
    )
    # Левый блок — что есть
    _rect(s, Inches(0.55), Inches(3.0), Inches(0.04), Inches(3.7), MINT)
    _txt(s, Inches(0.75), Inches(3.0), Inches(5.5), Inches(0.5),
         "Сделано", font="Cambria", size=20, bold=True, color=INK)
    _bullets(s, Inches(0.75), Inches(3.5), Inches(5.5), Inches(3.4),
             [
                 "Cortex-M3 ядро + 16 форматов Thumb + Thumb-2.",
                 "NVIC на 43 IRQ + SysTick.",
                 "TIM2, USART1, GPIO×5, RCC.",
                 "GDB RSP сервер на порту 3333.",
                 "11 наборов тестов, make test зелёный.",
                 "Пример прошивки и автогенератор stub-ELF.",
             ],
             color=INK, size=14, bullet_color=MINT)
    # Правый блок — команда (заполнить)
    _rect(s, Inches(6.85), Inches(3.0), Inches(0.04), Inches(3.7), AMBER)
    _txt(s, Inches(7.05), Inches(3.0), Inches(5.7), Inches(0.5),
         "Команда (вписать)", font="Cambria", size=20, bold=True, color=INK)
    _bullets(s, Inches(7.05), Inches(3.5), Inches(5.7), Inches(3.4),
             [
                 "Архитектор / ядро CPU — _____________________",
                 "Память + загрузчик — _____________________",
                 "NVIC + прерывания — _____________________",
                 "TIM2 + SysTick — _____________________",
                 "USART1 + CLI — _____________________",
                 "Тесты, сборка, GDB — _____________________",
             ],
             color=INK, size=14, bullet_color=AMBER)

    # Финал — снова cover
    s = _blank(prs)
    _bg(s, NAVY)
    _rect(s, Inches(0), Inches(0), Inches(0.10), SLIDE_H, AMBER)
    _txt(s, Inches(0.85), Inches(2.4), Inches(11), Inches(0.6),
         "// q&a", font="Consolas", size=14, color=AMBER_SOFT)
    _txt(s, Inches(0.85), Inches(3.0), Inches(11.6), Inches(2),
         "Спасибо. Готовы ответить на вопросы.",
         font="Cambria", size=44, bold=True, color=WHITE,
         line_spacing=1.05)
    _rect(s, Inches(0.85), Inches(5.0), Inches(0.6), Inches(0.04), AMBER)
    _txt(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(0.5),
         "STM32F103 Simulator · ИТМО · ПрВС 2026",
         font="Consolas", size=13, color=DIM_LIGHT)

    prs.save(str(OUT))
    print(f"Saved {OUT} ({len(prs.slides)} slides)")
    return OUT


if __name__ == "__main__":
    build()
