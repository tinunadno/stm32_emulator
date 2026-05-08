#!/usr/bin/env python3
"""
Защита проекта STM32-симулятора — финальная презентация.

Палитра: Charcoal + Amber (LED-вайб embedded-проектов).
Мотив: тонкая амбер-полоса слева + цифры в кружках.

Запуск:
    python3 presentation/build_defense_v2.py
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
HERE = Path(__file__).resolve().parent
OUT = HERE / "Защита_STM32_финал.pptx"

# --- Палитра (Charcoal Minimal + Amber accent) ---
DARK = RGBColor(0x1A, 0x1F, 0x2E)
DARK2 = RGBColor(0x26, 0x2D, 0x3F)
PAPER = RGBColor(0xF5, 0xF1, 0xE8)
PAPER_SOFT = RGBColor(0xEC, 0xE7, 0xDB)
AMBER = RGBColor(0xFF, 0xB3, 0x00)
AMBER_DARK = RGBColor(0xD6, 0x97, 0x00)
TEAL = RGBColor(0x2B, 0x7A, 0x78)
TEXT = RGBColor(0x1A, 0x1F, 0x2E)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
INK = RGBColor(0xFF, 0xFF, 0xFF)

# Шрифты
H_FONT = "Consolas"
B_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- helpers ----------

def load_image_generator():
    path = HERE / "build_full_md_deck.py"
    name = "_stm32_deck_images_v2"
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    assert spec.loader
    spec.loader.exec_module(mod)
    return mod.generate_all_images


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, rgb):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = rgb


def add_rect(slide, left, top, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background() if line is None else (
        setattr(sh.line, "color", None) or sh.line.color.rgb
    )
    if line is None:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_amber_strip(slide):
    """Тонкая амбер-полоса слева — сквозной мотив."""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    sh.fill.solid()
    sh.fill.fore_color.rgb = AMBER
    sh.line.fill.background()


def add_text(slide, left, top, w, h, text, *, font=B_FONT, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = text.split("\n")
    else:
        lines = list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
    return box


def add_circle_number(slide, left, top, n: int, diameter=Inches(0.65), color=AMBER, text_color=DARK):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = str(n)
    for r in p.runs:
        r.font.name = H_FONT
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = text_color


def add_amber_dot(slide, cx, cy, r=Inches(0.08)):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    sh.fill.solid()
    sh.fill.fore_color.rgb = AMBER
    sh.line.fill.background()


def add_kicker(slide, top_in: float, label: str):
    """Маленькая надстрочная подпись над заголовком (например «Раздел 2 · NVIC»)."""
    add_text(
        slide, Inches(0.55), Inches(top_in), Inches(11.5), Inches(0.3), label,
        font=H_FONT, size=11, bold=True, color=AMBER_DARK,
    )


def add_footer_page(slide, page: int, total: int, section: str = ""):
    add_text(
        slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.4),
        section,
        font=H_FONT, size=10, color=MUTED,
    )
    add_text(
        slide, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.4),
        f"{page:02d} / {total:02d}",
        font=H_FONT, size=10, color=MUTED, align=PP_ALIGN.RIGHT,
    )


# ---------- layouts ----------

def slide_cover(prs, title: str, subtitle: str, course: str):
    s = add_blank(prs)
    fill_bg(s, DARK)
    # амбер блок-логотип
    add_rect(s, Inches(0.55), Inches(0.55), Inches(0.55), Inches(0.55), AMBER)
    add_text(s, Inches(1.3), Inches(0.55), Inches(5.0), Inches(0.55),
             "STM32 SIMULATOR", font=H_FONT, size=14, bold=True, color=AMBER, anchor=MSO_ANCHOR.MIDDLE)

    # большие хайлайты в правом верхнем углу
    add_text(s, Inches(8.4), Inches(0.55), Inches(4.4), Inches(0.55),
             "ИТМО · 2026", font=H_FONT, size=14, color=PAPER, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # заголовок
    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(1.2),
             title, font=H_FONT, size=44, bold=True, color=PAPER)

    # подзаголовок
    add_text(s, Inches(0.7), Inches(3.55), Inches(12), Inches(1.2),
             subtitle, font=B_FONT, size=22, color=AMBER, italic=True, line_spacing=1.2)

    # горизонтальная амбер-черта
    add_rect(s, Inches(0.7), Inches(5.0), Inches(2.0), Inches(0.04), AMBER)

    # курс
    add_text(s, Inches(0.7), Inches(5.15), Inches(12), Inches(0.6),
             course, font=B_FONT, size=15, color=PAPER, line_spacing=1.3)

    # нижние «LED-точки» как мотив
    for i in range(4):
        add_amber_dot(s, Inches(0.7 + i * 0.25), Inches(7.0))


def slide_section_divider(prs, num: int, total: int, title: str, hint: str):
    s = add_blank(prs)
    fill_bg(s, DARK)
    add_amber_strip(s)
    add_text(s, Inches(0.7), Inches(2.0), Inches(8), Inches(0.6),
             f"РАЗДЕЛ {num:02d} / {total:02d}", font=H_FONT, size=16, bold=True, color=AMBER)
    add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(2.0),
             title, font=H_FONT, size=44, bold=True, color=PAPER, line_spacing=1.05)
    add_text(s, Inches(0.7), Inches(5.0), Inches(11), Inches(1.6),
             hint, font=B_FONT, size=20, color=PAPER_SOFT, italic=True, line_spacing=1.3)


def slide_one_idea(prs, kicker: str, title: str, body: str, page: int, total: int):
    """Большая мысль одним абзацем — для важных утверждений."""
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(1.3),
             title, font=H_FONT, size=34, bold=True, color=DARK, line_spacing=1.05)
    add_text(s, Inches(0.55), Inches(2.7), Inches(12.3), Inches(3.6),
             body, font=B_FONT, size=22, color=TEXT, line_spacing=1.35)
    add_footer_page(s, page, total, kicker)


def slide_stats(prs, kicker: str, title: str, stats: list[tuple[str, str]], note: str, page: int, total: int):
    """Большие цифры. stats = [(value, label), …] — до 4."""
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.9),
             title, font=H_FONT, size=30, bold=True, color=DARK)

    n = len(stats)
    block_w = Inches(11.5 / max(n, 1))
    for i, (val, label) in enumerate(stats):
        x = Inches(0.55) + Inches(11.5 / n) * i
        add_rect(s, x + Inches(0.05), Inches(2.2), block_w - Inches(0.1), Inches(3.0), DARK)
        add_text(s, x + Inches(0.1), Inches(2.35), block_w - Inches(0.2), Inches(1.8),
                 val, font=H_FONT, size=64, bold=True, color=AMBER, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), Inches(4.25), block_w - Inches(0.2), Inches(0.85),
                 label, font=B_FONT, size=14, color=PAPER, align=PP_ALIGN.CENTER, line_spacing=1.2)
    add_text(s, Inches(0.55), Inches(5.55), Inches(12.3), Inches(1.4),
             note, font=B_FONT, size=18, color=TEXT, italic=True, line_spacing=1.3)
    add_footer_page(s, page, total, kicker)


def slide_two_cols(prs, kicker: str, title: str, left_title: str, left_lines: list[str], right_title: str, right_lines: list[str], page: int, total: int):
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.9),
             title, font=H_FONT, size=28, bold=True, color=DARK)

    col_w = Inches(5.85)
    col_h = Inches(5.0)
    add_rect(s, Inches(0.55), Inches(2.05), col_w, col_h, PAPER_SOFT)
    add_rect(s, Inches(6.95), Inches(2.05), col_w, col_h, PAPER_SOFT)

    # левая колонка
    add_text(s, Inches(0.85), Inches(2.2), col_w - Inches(0.6), Inches(0.6),
             left_title, font=H_FONT, size=18, bold=True, color=TEAL)
    body = "\n".join("• " + s_ for s_ in left_lines)
    add_text(s, Inches(0.85), Inches(2.85), col_w - Inches(0.6), col_h - Inches(0.9),
             body, font=B_FONT, size=15, color=TEXT, line_spacing=1.35)

    # правая
    add_text(s, Inches(7.25), Inches(2.2), col_w - Inches(0.6), Inches(0.6),
             right_title, font=H_FONT, size=18, bold=True, color=AMBER_DARK)
    body = "\n".join("• " + s_ for s_ in right_lines)
    add_text(s, Inches(7.25), Inches(2.85), col_w - Inches(0.6), col_h - Inches(0.9),
             body, font=B_FONT, size=15, color=TEXT, line_spacing=1.35)
    add_footer_page(s, page, total, kicker)


def slide_picture(prs, kicker: str, title: str, png: Path, before: str, after: str, page: int, total: int, pic_h_in: float = 4.1):
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.7),
             title, font=H_FONT, size=26, bold=True, color=DARK)
    add_text(s, Inches(0.55), Inches(1.7), Inches(12.3), Inches(0.85),
             before, font=B_FONT, size=15, color=TEXT, italic=True, line_spacing=1.25)

    pic_h_in = min(pic_h_in, 3.4)  # запас под нижнюю подпись и футер
    if png.is_file():
        s.shapes.add_picture(str(png), Inches(1.0), Inches(2.55), height=Inches(pic_h_in))

    after_top = 2.55 + pic_h_in + 0.10
    after_h = max(0.5, 6.95 - after_top)
    add_text(s, Inches(0.55), Inches(after_top), Inches(12.3), Inches(after_h),
             after, font=B_FONT, size=13, color=TEXT, line_spacing=1.25)
    add_footer_page(s, page, total, kicker)


def slide_icon_grid(prs, kicker: str, title: str, items: list[tuple[str, str, str]], page: int, total: int):
    """3×2 сетка карточек: (метка-кружок, заголовок, описание)."""
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.9),
             title, font=H_FONT, size=28, bold=True, color=DARK)

    cols, rows = 3, 2
    cw = Inches(4.05)
    ch = Inches(2.45)
    gx = Inches(0.20)
    gy = Inches(0.20)
    x0 = Inches(0.55)
    y0 = Inches(2.05)
    for i, (label, htxt, body) in enumerate(items[:cols * rows]):
        c = i % cols
        r = i // cols
        x = x0 + (cw + gx) * c
        y = y0 + (ch + gy) * r
        add_rect(s, x, y, cw, ch, PAPER_SOFT)
        # маленький амбер кружок
        add_circle_number_text(s, x + Inches(0.25), y + Inches(0.25), label)
        add_text(s, x + Inches(1.1), y + Inches(0.2), cw - Inches(1.2), Inches(0.6),
                 htxt, font=H_FONT, size=14, bold=True, color=DARK)
        add_text(s, x + Inches(0.25), y + Inches(0.95), cw - Inches(0.4), ch - Inches(1.05),
                 body, font=B_FONT, size=12, color=TEXT, line_spacing=1.3)
    add_footer_page(s, page, total, kicker)


def add_circle_number_text(slide, left, top, label: str, diameter=Inches(0.7), color=AMBER, text_color=DARK):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    for r in p.runs:
        r.font.name = H_FONT
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = text_color


def slide_code(prs, kicker: str, title: str, before: str, code: str, after: str, page: int, total: int):
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.7),
             title, font=H_FONT, size=26, bold=True, color=DARK)
    add_text(s, Inches(0.55), Inches(1.7), Inches(12.3), Inches(0.7),
             before, font=B_FONT, size=14, color=TEXT, italic=True, line_spacing=1.3)
    # код-блок
    add_rect(s, Inches(0.55), Inches(2.45), Inches(12.2), Inches(3.5), DARK)
    add_text(s, Inches(0.75), Inches(2.6), Inches(11.8), Inches(3.2),
             code, font="Consolas", size=14, color=PAPER, line_spacing=1.3)
    add_text(s, Inches(0.55), Inches(6.05), Inches(12.3), Inches(0.85),
             after, font=B_FONT, size=14, color=TEXT, line_spacing=1.3)
    add_footer_page(s, page, total, kicker)


def slide_callout_quote(prs, kicker: str, quote: str, attribution: str, page: int, total: int):
    s = add_blank(prs)
    fill_bg(s, DARK)
    add_text(s, Inches(0.7), Inches(0.55), Inches(12), Inches(0.5),
             kicker, font=H_FONT, size=14, bold=True, color=AMBER)
    add_rect(s, Inches(0.7), Inches(2.0), Inches(0.06), Inches(3.7), AMBER)
    add_text(s, Inches(1.0), Inches(2.0), Inches(11.5), Inches(3.7),
             quote, font=H_FONT, size=30, bold=True, color=PAPER, line_spacing=1.2)
    add_text(s, Inches(1.0), Inches(5.85), Inches(11.5), Inches(0.8),
             attribution, font=B_FONT, size=16, color=AMBER, italic=True)
    add_text(s, Inches(11.5), Inches(7.05), Inches(1.4), Inches(0.4),
             f"{page:02d} / {total:02d}",
             font=H_FONT, size=10, color=PAPER_SOFT, align=PP_ALIGN.RIGHT)


def slide_table(prs, kicker: str, title: str, header: list[str], rows: list[list[str]], note: str, page: int, total: int):
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.9),
             title, font=H_FONT, size=28, bold=True, color=DARK)

    # заголовок таблицы
    n = len(header)
    total_w = Inches(12.2)
    col_w = Emu(int(total_w.emu / n))
    y = Inches(2.05)
    add_rect(s, Inches(0.55), y, total_w, Inches(0.55), DARK)
    for i, h in enumerate(header):
        add_text(s, Inches(0.55) + col_w * i + Inches(0.15), y + Inches(0.07), col_w - Inches(0.2), Inches(0.5),
                 h, font=H_FONT, size=14, bold=True, color=AMBER)

    # строки
    row_h = Inches(0.55)
    for ri, row in enumerate(rows):
        ry = y + Inches(0.55) + row_h * ri
        bg = PAPER if ri % 2 == 0 else PAPER_SOFT
        add_rect(s, Inches(0.55), ry, total_w, row_h, bg)
        for ci, cell in enumerate(row):
            add_text(s, Inches(0.55) + col_w * ci + Inches(0.15), ry + Inches(0.1), col_w - Inches(0.2), row_h - Inches(0.15),
                     cell, font=B_FONT, size=13, color=TEXT, line_spacing=1.2)

    add_text(s, Inches(0.55), Inches(6.45), Inches(12.3), Inches(0.55),
             note, font=B_FONT, size=13, color=MUTED, italic=True)
    add_footer_page(s, page, total, kicker)


def slide_timeline(prs, kicker: str, title: str, steps: list[tuple[str, str]], note: str, page: int, total: int):
    """Шаги пайплайна: квадраты с номерами и стрелками."""
    s = add_blank(prs)
    fill_bg(s, PAPER)
    add_amber_strip(s)
    add_kicker(s, 0.55, kicker)
    add_text(s, Inches(0.55), Inches(0.95), Inches(12), Inches(0.9),
             title, font=H_FONT, size=28, bold=True, color=DARK)

    n = len(steps)
    total_w = Inches(12.2)
    box_w = Emu(int((total_w.emu - Inches(0.4 * (n - 1)).emu) / n))
    box_h = Inches(2.6)
    y = Inches(2.5)
    x0 = Inches(0.55)
    for i, (head, body) in enumerate(steps):
        x = x0 + (box_w + Inches(0.4)) * i
        add_rect(s, x, y, box_w, box_h, DARK)
        add_text(s, x + Inches(0.2), y + Inches(0.2), box_w - Inches(0.4), Inches(0.5),
                 f"0{i+1}", font=H_FONT, size=12, bold=True, color=AMBER)
        add_text(s, x + Inches(0.2), y + Inches(0.65), box_w - Inches(0.4), Inches(0.7),
                 head, font=H_FONT, size=18, bold=True, color=PAPER, line_spacing=1.1)
        add_text(s, x + Inches(0.2), y + Inches(1.45), box_w - Inches(0.4), box_h - Inches(1.6),
                 body, font=B_FONT, size=12, color=PAPER_SOFT, line_spacing=1.3)
        if i < n - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, x + box_w + Inches(0.05), y + Inches(1.05), Inches(0.3), Inches(0.5))
            arr.fill.solid()
            arr.fill.fore_color.rgb = AMBER
            arr.line.fill.background()
    add_text(s, Inches(0.55), Inches(5.6), Inches(12.3), Inches(1.0),
             note, font=B_FONT, size=15, color=TEXT, italic=True, line_spacing=1.3)
    add_footer_page(s, page, total, kicker)


def slide_closing(prs, title: str, sub: str):
    s = add_blank(prs)
    fill_bg(s, DARK)
    add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.2),
             title, font=H_FONT, size=64, bold=True, color=AMBER)
    add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(2),
             sub, font=B_FONT, size=22, color=PAPER, italic=True, line_spacing=1.3)
    for i in range(6):
        add_amber_dot(s, Inches(0.7 + i * 0.25), Inches(6.7))


# ---------- build ----------

def build():
    img = load_image_generator()()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Сначала собираем все вызовы — потом вычислим total и подставим page
    builders = []

    builders.append(lambda p, t: slide_cover(
        prs,
        "Симулятор STM32F103C8T6",
        "программная модель Cortex-M3 c периферией, тестами и удалённой отладкой",
        "Курс «Программирование встроенных систем» · ИТМО · защита группового проекта",
    ))

    # ============== I. ВВЕДЕНИЕ ==============
    builders.append(lambda p, t: slide_section_divider(prs, 1, 6, "Что это и зачем", "Контекст проекта в одну минуту"))

    builders.append(lambda p, t: slide_one_idea(
        prs,
        "Идея",
        "Программа, которая ведёт себя как микроконтроллер.",
        "Мы написали на чистом C11 модель «Blue Pill» — STM32F103C8T6. "
        "Загружаем готовую прошивку, симулируем процессор, память и периферию, "
        "получаем такие же эффекты, как на реальной плате — и всё это без железа.",
        p, t,
    ))

    builders.append(lambda p, t: slide_stats(
        prs,
        "Числами",
        "Из чего состоит проект",
        [
            ("≈ 7 400", "строк C-кода\nв src/"),
            ("30", "файлов .c\n+ 20 .h"),
            ("11", "тестовых сьютов\n+ интеграция"),
            ("~80", "Thumb / Thumb-2\nинструкций"),
        ],
        "Цифры собраны автоматически по содержимому каталога src/. "
        "Это не «кодовая фабрика» — это плотный учебный проект, в котором каждая папка делает одно дело.",
        p, t,
    ))

    builders.append(lambda p, t: slide_one_idea(
        prs,
        "Зачем студенту симулятор",
        "Видеть процессор изнутри без риска что-нибудь спалить.",
        "Симулятор даёт пошаговую трассировку, доступ к каждому регистру, "
        "и привычный GDB как на железе. Прошивку можно собирать своим компилятором — "
        "и сразу запускать. Это удобно для лабораторных и для отладки алгоритмов на реальной модели Cortex-M3.",
        p, t,
    ))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "Что было в задании",
        "Курсовое ТЗ → реализация",
        "Что требовалось",
        [
            "ARM Cortex-M3 (Thumb-2)",
            "Подсистема памяти Flash + SRAM",
            "NVIC c приоритетами",
            "Таймер, UART, MMIO",
            "Загрузка прошивки и тесты",
        ],
        "Что мы сделали сверх",
        [
            "Thumb-2: BL/MOVW/MOVT/UMULL/SMULL",
            "MMIO-NVIC в стиле CMSIS",
            "GDB Remote Serial Protocol",
            "SysTick + RCC + GPIO",
            "Профилировщик и UART SVG",
        ],
        p, t,
    ))

    # ============== II. АРХИТЕКТУРА ==============
    builders.append(lambda p, t: slide_section_divider(prs, 2, 6, "Архитектура", "Один поток, чёткие границы между модулями"))

    builders.append(lambda p, t: slide_picture(
        prs,
        "Картинка 1 · Архитектура",
        "Как соединены ядро, шина, память и периферия",
        img["architecture"],
        "Слева — процессор Cortex-M3, справа — NVIC. Между ними — шина (bus), которая по адресу решает, "
        "к какой периферии или памяти идёт обращение. Так ядро не «знает» про устройства напрямую.",
        "Все стрелки на этом слайде соответствуют реальным вызовам в коде: bus_read/bus_write, "
        "nvic_set_pending, register_region. Источник: docs/architecture.svg.",
        p, t, pic_h_in=3.7,
    ))

    builders.append(lambda p, t: slide_picture(
        prs,
        "Картинка 2 · Каталог src/",
        "Кто за что отвечает в коде",
        img["modules"],
        "Каждая папка — отдельная зона ответственности. Это упрощает работу команды: "
        "разработчик ядра не зависит от UART-разработчика, тесты модулей независимы.",
        "core — процессор; memory — Flash/SRAM; bus — маршрутизация; nvic — прерывания; "
        "peripherals — устройства; simulator — главный цикл; ui и gdb_stub — интерфейсы.",
        p, t, pic_h_in=3.5,
    ))

    builders.append(lambda p, t: slide_picture(
        prs,
        "Картинка 3 · Карта памяти",
        "Адресное пространство такое же, как в datasheet STM32F103",
        img["memory_map"],
        "Прошивка обращается к регистрам по тем же адресам, что и на кристалле: 0x4000_0000 — TIM2, "
        "0x4001_3800 — USART1. Поэтому код, собранный под реальный чип, работает в симуляторе без правок.",
        "Особенность: Flash доступен как 0x0000_0000 (alias для векторной таблицы при сбросе) "
        "и как 0x0800_0000 — это поведение настоящего STM32.",
        p, t, pic_h_in=3.0,
    ))

    builders.append(lambda p, t: slide_picture(
        prs,
        "Картинка 4 · Один такт симулятора",
        "Жёсткий порядок шагов даёт детерминизм",
        img["sim_cycle"],
        "Сначала тикают периферии (таймеры могут поднять IRQ), потом исполняется одна инструкция, "
        "потом NVIC решает, нужно ли войти в обработчик, и в конце проверяются breakpoints.",
        "Этот порядок жёстко зафиксирован в simulator_step и не меняется. Тесты опираются на эту "
        "повторяемость — поэтому они не «мигают» от запуска к запуску.",
        p, t, pic_h_in=2.4,
    ))

    # ============== III. ЯДРО И ПЕРИФЕРИЯ ==============
    builders.append(lambda p, t: slide_section_divider(prs, 3, 6, "Ядро и периферия", "Что именно умеет наш процессор"))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "Cortex-M3 core/",
        "Процессор: чем он умеет «думать»",
        "Состояние",
        [
            "Регистры R0–R15, xPSR (N/Z/C/V)",
            "Флаги: thumb_mode, interruptible",
            "Счётчик cycles (uint64)",
            "Текущий IRQ — для вложенности",
        ],
        "Декодер",
        [
            "Таблица 16-бит: маска + паттерн",
            "Отдельный путь для Thumb-2 32-bit",
            "BL, MOVW/MOVT, UMULL, SMULL, MUL",
            "Условные ветвления по флагам",
        ],
        p, t,
    ))

    builders.append(lambda p, t: slide_code(
        prs,
        "core/ · core_step",
        "Один такт работы ядра",
        "Так выглядит сердце симулятора: декод инструкции, исполнение, обновление регистров.",
        "Status core_step(Core* core)\n"
        "{\n"
        "    uint16_t instr = bus_read(core->bus, core->state.r[15], 2);\n"
        "    if ((instr & 0xE000) == 0xE000 && (instr & 0x1800)) {\n"
        "        // Thumb-2: BL, MOVW/MOVT, UMULL, …\n"
        "        execute_32bit(core, instr, hw2);\n"
        "    } else {\n"
        "        // Линейный поиск по таблице 16-бит инструкций\n"
        "        for (entry = instr_table; entry->handler; ++entry)\n"
        "            if ((instr & entry->mask) == entry->pattern) break;\n"
        "        entry->handler(core, instr);\n"
        "    }\n"
        "    CYCLES++;\n"
        "    if (core->state.interruptible)\n"
        "        if (nvic_get_pending_irq(...)) enter_exception(...);\n"
        "}",
        "Таблица отсортирована: более специфичные паттерны сверху. Это типичный приём "
        "для интерпретаторов ARM — простой, расширяемый, легко покрывается тестами.",
        p, t,
    ))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "nvic/ · контроллер прерываний",
        "Кто и когда прерывает наш CPU",
        "Что хранится",
        [
            "До 43 линий IRQ + SysTick",
            "pending / active / enabled / priority",
            "current_priority — кто сейчас",
            "Меньше число — выше приоритет",
        ],
        "Что делает",
        [
            "set_pending — периферия просит",
            "get_pending_irq — выбор по приоритету",
            "acknowledge / complete — вход / выход",
            "MMIO в стиле CMSIS (ISER, ICER, IPR)",
        ],
        p, t,
    ))

    builders.append(lambda p, t: slide_one_idea(
        prs,
        "nvic/ · MMIO в стиле CMSIS",
        "Прошивка работает с прерываниями как с обычной памятью.",
        "В файле nvic_bus.c регистры NVIC отображены на адреса по правилам ARM CMSIS: "
        "ISER/ICER (включить/выключить IRQ), ISPR/ICPR (поднять/снять pending), IPR (приоритеты). "
        "Это значит, что HAL и stm32f1xx.h видят наш симулятор так же, как настоящий чип.",
        p, t,
    ))

    builders.append(lambda p, t: slide_icon_grid(
        prs,
        "Периферия",
        "Шесть устройств — каждое со своим vtable",
        [
            ("01", "TIM2", "Базовый таймер с PSC и ARR. UIF в SR; событие в event_queue даёт IRQ 28."),
            ("02", "USART1", "TX и RX, прерывания TXEIE / RXNEIE → IRQ 37; SVG-диаграмма для отладки."),
            ("03", "SysTick", "Системный таймер CPU. TICKINT → исключение 15, обходит NVIC IRQ-линии."),
            ("04", "GPIO A/B/C", "CRL/CRH/IDR/ODR/BSRR. Можно слать пины снаружи через gpio_set_pin."),
            ("05", "RCC", "Упрощённо: ready-биты следуют за enable, SWS = SW. HAL не виснет."),
            ("06", "NVIC bus", "Регистры контроллера прерываний прямо как MMIO — для драйверов CMSIS."),
        ],
        p, t,
    ))

    builders.append(lambda p, t: slide_picture(
        prs,
        "Картинка 5 · Диаграмма UART",
        "Когда какой байт ушёл по линии TX",
        img["uart"],
        "Симулятор умеет выводить SVG-диаграмму: каждый «столбик» — это байт на TX в момент "
        "симуляционного такта. Очень удобно показывать на защите, что прошивка реально шевелится.",
        "Эту диаграмму можно сохранить из CLI командой diagram-save или запросить из GDB через "
        "monitor uart-svg. Источник графики: uart_diagram.svg.",
        p, t, pic_h_in=2.6,
    ))

    # ============== IV. ОТЛАДКА И ТЕСТЫ ==============
    builders.append(lambda p, t: slide_section_divider(prs, 4, 6, "Отладка и тесты", "Чем доказываем, что работает"))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "gdb_stub/ · удалённая отладка",
        "Привычный GDB подключается к симулятору по TCP",
        "Поддержанные пакеты RSP",
        [
            "? · g · G · p · P  — регистры",
            "m · M  — память (через шину)",
            "c · s  — continue / step",
            "Z · z  — программные breakpoints",
            "qXfer:features:read  — target.xml",
        ],
        "Monitor-команды",
        [
            "monitor halt / reset / reset halt",
            "monitor uart-diagram / uart-clear",
            "monitor uart-svg",
            "qSupported, H, T, D, k — без боли",
            "Как у настоящего gdbserver",
        ],
        p, t,
    ))

    builders.append(lambda p, t: slide_table(
        prs,
        "ui/ · CLI-команды",
        "Самые нужные команды интерактивного режима",
        ["Команда", "Что делает"],
        [
            ["help", "список всех команд"],
            ["load <путь>", "загрузить .bin во Flash"],
            ["step [N] · run · stop", "выполнение и пауза"],
            ["reg · mem <addr>", "состояние процессора"],
            ["break <addr> · delete", "точки останова"],
            ["uart <ch> · diagram-save", "ввод и SVG-вывод UART"],
            ["profile · quit", "профилировщик и выход"],
        ],
        "Парсер написан вручную — простой strcmp по таблице команд. Это видно начинающему и "
        "легко расширяется новыми командами без правки старого кода.",
        p, t,
    ))

    builders.append(lambda p, t: slide_stats(
        prs,
        "Тесты",
        "Что проверяет make test",
        [
            ("11", "сьютов\nпо модулям"),
            ("3", "интеграционных\nсценария"),
            ("100%", "ключевых\nмодулей"),
            ("ASSERT*", "макросы\nс счётчиком"),
        ],
        "Каркас минимальный: ASSERT_EQ при провале печатает файл, строку и условие, "
        "увеличивает счётчик и выходит из теста. RUN_TEST печатает PASS / FAIL и в конце — итог.",
        p, t,
    ))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "Тесты · что проверяем",
        "Юнит и интеграция бок о бок",
        "Юнит-тесты",
        [
            "core: MOV, ADD, CMP+B, PUSH/POP, MUL, BL+BX",
            "memory + bus: alias 0x0/0x08000000",
            "nvic: приоритеты и вложенность",
            "timer/systick: переполнение и IRQ",
            "uart/gpio/rcc: регистры и поведение",
        ],
        "Интеграция (test_integration.c)",
        [
            "TIM2 → NVIC → handler → BX LR",
            "Breakpoint останавливает simulator_run",
            "UART через bus → callback пользователя",
            "Соединяем CPU + периферию",
            "Это и есть «прошивка работает»",
        ],
        p, t,
    ))

    # ============== V. КАК ПОЛЬЗОВАТЬСЯ ==============
    builders.append(lambda p, t: slide_section_divider(prs, 5, 6, "Как пользоваться", "Сборка, запуск, демо"))

    builders.append(lambda p, t: slide_code(
        prs,
        "Сборка хоста",
        "Запуск за минуту",
        "Симулятор собирается как обычная C-программа на Linux или macOS.",
        "$ cd src\n"
        "$ make                 # CC=clang, -std=c11, -Wall -Wextra -Wpedantic\n"
        "$ ./stm32sim           # интерактивный CLI\n"
        "$ ./stm32sim fw.bin    # сразу с прошивкой\n"
        "$ ./stm32sim fw.bin --gdb 4444   # GDB-сервер на порту 4444\n"
        "$ make test            # собрать и запустить test_runner",
        "В Makefile нет санитайзеров — это сознательно: тестам важнее воспроизводимость, "
        "чем диагностика памяти. ASAN можно включить руками через CFLAGS.",
        p, t,
    ))

    builders.append(lambda p, t: slide_code(
        prs,
        "Сборка прошивки",
        "Своя прошивка для симулятора",
        "Пример из examples/ собирается обычным arm-none-eabi-gcc — как для настоящего STM32.",
        "$ cd examples\n"
        "$ make                              # → firmware.elf, firmware.bin\n"
        "# CFLAGS = -mcpu=cortex-m3 -mthumb -nostdlib -ffreestanding -O0 -g\n"
        "# LDFLAGS = -T link.ld -nostdlib -nostartfiles\n"
        "$ ../src/stm32sim firmware.bin      # запуск в симуляторе\n"
        "# или без gcc:\n"
        "$ python3 gen_firmware.py           # мини-ассемблер Thumb на Python",
        "В examples/ лежит реалистичная прошивка: vector table, RCC HSI→PLL, GPIOA PA5, USART1, "
        "SysTick и TIM2 c обработчиками. Хороший пример «как должно выглядеть».",
        p, t,
    ))

    builders.append(lambda p, t: slide_timeline(
        prs,
        "Демо на защите · 2 минуты",
        "Пошаговый сценарий, который не подведёт",
        [
            ("Сборка", "make в src/ и examples/. Покажите zero warnings."),
            ("CLI", "load fw.bin, step 5, reg, mem 0x20000000."),
            ("UART", "diagram-save: SVG c байтами на линии."),
            ("GDB", "stm32sim --gdb · target remote · breakpoint."),
        ],
        "Если что-то не собирается на чужой машине — заранее запускаем make test и показываем зелёные галочки. "
        "Это сильнее любых слов о «корректности».",
        p, t,
    ))

    # ============== VI. КОМАНДА И ИТОГ ==============
    builders.append(lambda p, t: slide_section_divider(prs, 6, 6, "Команда и итог", "Кто это сделал и что дальше"))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "Распределение ролей",
        "По шести областям — как в задании",
        "Кто что делал",
        [
            "1. Архитектор · ядро CPU",
            "2. Память и загрузка прошивки",
            "3. NVIC и обработчики",
            "4. Таймеры (TIM2, SysTick)",
            "5. UART и GPIO/RCC",
            "6. Тесты, сборка, документация",
        ],
        "Этапы по семестру",
        [
            "проект архитектуры",
            "ядро + память + шина",
            "NVIC + интеграция",
            "периферия (TIM, UART, GPIO)",
            "GDB + CLI + профилировщик",
            "тесты, отчёт, защита",
        ],
        p, t,
    ))

    builders.append(lambda p, t: slide_picture(
        prs,
        "Что дальше",
        "Идея онлайн-лаборатории (только в документе docs/gateway.md)",
        img["gateway"],
        "В репозитории есть проектная заметка: распределённая система с очередью KeyDB, "
        "воркерами на Go и SSE для браузера. Это план развития, а не текущий код.",
        "Сейчас результат — локальный симулятор с тестами и GDB. Распределённую оболочку "
        "можно собирать поверх него уже без переделки ядра.",
        p, t, pic_h_in=2.8,
    ))

    builders.append(lambda p, t: slide_two_cols(
        prs,
        "Честный статус",
        "Что есть и чего пока нет — без приукрашивания",
        "Сделано",
        [
            "C11 без зависимостей",
            "Cortex-M3 + Thumb-2 ядро",
            "Шина, MMIO, прерывания, SysTick",
            "GDB RSP + CLI + профилировщик",
            "11 тестов + интеграция",
        ],
        "Не сделано (и это нормально)",
        [
            "Полный набор периферии чипа",
            "DMA, SPI/I2C, ADC",
            "CI в репозитории",
            "Онлайн-сервис из gateway.md",
            "Идеальное соответствие тактам кремния",
        ],
        p, t,
    ))

    builders.append(lambda p, t: slide_callout_quote(
        prs,
        "Главная мысль",
        "Это не «эмуляция железа в реальном времени» — это детерминированная учебная модель,"
        " где каждое действие можно объяснить и повторить.",
        "— почему наш симулятор полезен на занятиях",
        p, t,
    ))

    builders.append(lambda p, t: slide_one_idea(
        prs,
        "Команда (заполнить перед защитой)",
        "Шесть человек — шесть зон ответственности.",
        "ФИО · группа · роль (1):  ____________________________________________\n"
        "ФИО · группа · роль (2):  ____________________________________________\n"
        "ФИО · группа · роль (3):  ____________________________________________\n"
        "ФИО · группа · роль (4):  ____________________________________________\n"
        "ФИО · группа · роль (5):  ____________________________________________\n"
        "ФИО · группа · роль (6):  ____________________________________________",
        p, t,
    ))

    builders.append(lambda p, t: slide_closing(prs, "Спасибо.", "Готовы ответить на вопросы и показать симулятор в работе."))

    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        fn(i, total)

    prs.save(str(OUT))
    print(f"Saved {OUT} ({total} slides)")
    return OUT


if __name__ == "__main__":
    build()
